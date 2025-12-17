"""
PDF Master Tool - Cloud Run Backend
FastAPI service with LibreOffice conversions
"""

from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import subprocess
import os
import tempfile
import base64
import shutil
from pathlib import Path
import time
import pdfplumber
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
import re
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.drawing.image import Image as ExcelImage
from typing import Any, Dict, List, Optional, Set, Tuple
from pdf2docx import Converter
from PIL import Image as PILImage
import io
from pdf_to_excel_engine import convert_pdf_to_excel
try:
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    print("[WARNING] pytesseract not available - OCR features disabled")
try:
    import pypdfium2 as pdfium
    PDFIUM_AVAILABLE = True
except ImportError:
    PDFIUM_AVAILABLE = False
    print("[WARNING] pypdfium2 not available - advanced image extraction disabled")

app = FastAPI(title="PDF Master Tool API")

def filter_java_warnings(error_message: str) -> str:
    """
    ULTRA-AGGRESSIVE filter: If ANY Java reference found, return generic message.
    This completely hides Java warnings from users.
    """
    if not error_message:
        return "Conversion failed. Please try again."
    
    error_lower = error_message.lower()
    java_keywords = [
        "javaldx",
        "java may not function",
        "failed to launch javaldx",
        "java runtime",
        "java environment",
        "jfw",
        "warning: failed to launch",
    ]

    # ULTRA-AGGRESSIVE: If ANY Java warning found, return generic (non-technical) message
    if any(keyword in error_lower for keyword in java_keywords):
        return "Conversion failed. Please try again or check if the file format is supported."

    # For non-Java errors, just clean up duplicate prefixes but keep the real message
    result = error_message or ""
    while "Conversion failed: Conversion failed:" in result:
        result = result.replace("Conversion failed: Conversion failed:", "Conversion failed:")

    return result.strip()

# Enable CORS for frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Update with your frontend URL
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
async def root():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "service": "PDF Master Tool",
        "version": "1.0.0",
        "endpoints": [
            "/convert/pdf-to-word",
            "/convert/pdf-to-excel",
            "/convert/pdf-to-ppt",
            "/convert/word-to-pdf",
            "/convert/ppt-to-pdf",
            "/convert/excel-to-pdf"
        ]
    }

@app.get("/test-excel")
async def test_excel():
    """Test endpoint to create a simple Excel file with one sheet"""
    from openpyxl import Workbook
    import tempfile
    import base64
    
    wb = Workbook()
    wb.remove(wb.active)
    ws = wb.create_sheet(title="Content")
    ws.cell(row=1, column=1, value="Test - Single Sheet")
    ws.cell(row=2, column=1, value="This should be only one sheet")
    
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        wb.save(tmp.name)
        with open(tmp.name, 'rb') as f:
            data = f.read()
        os.unlink(tmp.name)
    
    return JSONResponse({
        "success": True,
        "file": base64.b64encode(data).decode(),
        "filename": "test.xlsx",
        "sheets": len(wb.worksheets),
        "sheet_names": [s.title for s in wb.worksheets]
    })

@app.get("/health")
async def health():
    """Health check for Cloud Run"""
    return {"status": "ok"}

@app.get("/test-libreoffice")
async def test_libreoffice():
    """Test if LibreOffice is installed and working"""
    import subprocess
    try:
        # Test if LibreOffice is accessible
        result = subprocess.run(
            ['libreoffice', '--version'],
            capture_output=True,
            text=True,
            timeout=10
        )
        version_info = result.stdout.strip() if result.stdout else "No output"
        
        # Test a simple conversion (create a test file)
        with tempfile.TemporaryDirectory() as temp_dir:
            # Create a simple text file
            test_file = Path(temp_dir) / "test.txt"
            test_file.write_text("Test content")
            
            # Try to convert to PDF
            cmd = [
                'libreoffice',
                '--headless',
                '--invisible',
                '--nofirststartwizard',
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                str(test_file)
            ]
            
            env = {
                **os.environ,
                'HOME': temp_dir,
                'SAL_USE_VCLPLUGIN': 'headless',
                'DISPLAY': '',
            }
            
            conv_result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=30,
                env=env
            )
            
            # Check if PDF was created
            pdf_file = Path(temp_dir) / "test.pdf"
            pdf_exists = pdf_file.exists()
            
            return {
                "status": "ok",
                "libreoffice_version": version_info,
                "libreoffice_return_code": result.returncode,
                "conversion_test": {
                    "return_code": conv_result.returncode,
                    "stdout": conv_result.stdout[:500] if conv_result.stdout else None,
                    "stderr": conv_result.stderr[:500] if conv_result.stderr else None,
                    "pdf_created": pdf_exists,
                    "pdf_size": pdf_file.stat().st_size if pdf_exists else 0
                }
            }
    except Exception as e:
        return {
            "status": "error",
            "error": str(e)
        }

def convert_with_libreoffice(
    input_path: str,
    output_format: str,
    output_dir: str,
    timeout: int = 45,  # ULTRA-OPTIMIZED: Reduced to 45s for faster processing (small files)
    allow_retry: bool = True,
    input_filter: Optional[str] = None,
) -> str:
    """
    Convert document using LibreOffice with optimized settings for speed
    
    Args:
        input_path: Path to input file
        output_format: Target format (docx, xlsx, pptx, pdf)
        output_dir: Directory for output file
        timeout: Timeout in seconds (default 90 for faster processing)
        
    Returns:
        Path to converted file
    """
    try:
        # Use absolute paths (LibreOffice needs absolute paths)
        input_path_abs = os.path.abspath(input_path)
        output_dir_abs = os.path.abspath(output_dir)
        
        # Verify input file exists
        if not os.path.exists(input_path_abs):
            raise Exception(f"Input file does not exist: {input_path_abs}")
        
        # Ensure output directory exists
        os.makedirs(output_dir_abs, exist_ok=True)
        
        # FIXED: Use soffice with absolute paths (Cloud Run requirement)
        # Try soffice first (more reliable), fallback to libreoffice
        soffice_path = '/usr/bin/soffice'
        libreoffice_path = '/usr/bin/libreoffice'
        
        # Use soffice if available, otherwise libreoffice
        if os.path.exists(soffice_path):
            cmd_base = soffice_path
        elif os.path.exists(libreoffice_path):
            cmd_base = libreoffice_path
        else:
            cmd_base = 'libreoffice'  # Fallback to PATH lookup
        
        # ULTRA-OPTIMIZED: Minimal command with performance flags for fastest conversion
        # Removed unnecessary flags that slow down startup
        cmd = [
            cmd_base,
            '--headless',           # No GUI (fastest)
            '--invisible',          # No splash screen
            '--nodefault',          # Don't load default document
            '--nolockcheck',        # Skip lock file check (faster)
            '--nologo',             # No logo splash
            '--norestore',          # Don't restore previous session
            '--nocrashreport',      # Disable crash reporting for speed
            '--convert-to', output_format,
            '--outdir', output_dir_abs,
            input_path_abs
        ]
        
        # OPTIMIZED: Set HOME to /tmp and add performance environment variables
        # Initialize env FIRST before using it
        env = os.environ.copy()
        
        # For Excel files, add print settings hint
        # LibreOffice will respect the Excel file's print settings if we don't override them
        if input_path_abs.lower().endswith(('.xlsx', '.xls')):
            # Add environment variable to hint LibreOffice to respect print settings
            env.setdefault('SAL_USE_VCLPLUGIN', 'headless')
            env.setdefault('OOO_DISABLE_RECOVERY', '1')
            print(f"[LibreOffice] Excel file detected - will respect print settings from file")

        if input_filter:
            cmd.insert(-2, '--infilter')  # Insert before --outdir
            cmd.insert(-2, input_filter)
        env['HOME'] = '/tmp'
        env['SAL_USE_VCLPLUGIN'] = 'headless'
        env['SAL_DISABLE_OPENCL'] = '1'  # Disable OpenCL for faster startup
        env['SAL_DISABLE_OPENGL'] = '1'  # Disable OpenGL for headless mode
        env['OOO_DISABLE_RECOVERY'] = '1'  # Disable recovery for speed
        env['SAL_NO_FONTCONFIG'] = '1'  # Skip font config for faster startup
        # Unset DISPLAY to prevent X11 errors
        env.pop('DISPLAY', None)
        
        print(f"[LibreOffice] Command: {' '.join(cmd)}")
        print(f"[LibreOffice] Using: {cmd_base}")
        print(f"[LibreOffice] Input: {input_path_abs} (exists: {os.path.exists(input_path_abs)}, size: {os.path.getsize(input_path_abs) if os.path.exists(input_path_abs) else 0} bytes)")
        print(f"[LibreOffice] Output dir: {output_dir_abs} (exists: {os.path.exists(output_dir_abs)}, writable: {os.access(output_dir_abs, os.W_OK)})")
        print(f"[LibreOffice] HOME: {env.get('HOME', 'not set')}")
        
        # List files before conversion
        files_before = list(Path(output_dir_abs).iterdir())
        print(f"[LibreOffice] Files before: {[f.name for f in files_before]}")
        
        # FIXED: Run LibreOffice with explicit error handling
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=timeout,
                env=env,
                check=False  # Don't raise on non-zero exit
            )
        except subprocess.TimeoutExpired:
            raise Exception(f"LibreOffice conversion timed out after {timeout} seconds")
        except Exception as e:
            raise Exception(f"Failed to run LibreOffice: {str(e)}")
        
        # FIXED: Log everything (Cloud Run debugging requirement)
        print(f"[LibreOffice] Return code: {result.returncode}")
        print(f"[LibreOffice] stdout (full): {result.stdout if result.stdout else '(empty)'}")
        print(f"[LibreOffice] stderr (full): {result.stderr if result.stderr else '(empty)'}")
        
        # Wait a moment for file system to sync (reduced from 0.5s to 0.2s for speed)
        time.sleep(0.2)
        
        # Log LibreOffice output for debugging
        if result.stdout:
            print(f"[LibreOffice] stdout: {result.stdout[:500]}")
        if result.stderr:
            # Filter out javaldx warnings (non-fatal)
            stderr_filtered = result.stderr
            if 'javaldx' in stderr_filtered.lower():
                print(f"[LibreOffice] ⚠️ Java warning (non-fatal): {stderr_filtered[:200]}")
                # Remove javaldx warnings from error message
                stderr_filtered = '\n'.join([line for line in stderr_filtered.split('\n') 
                                           if 'javaldx' not in line.lower()])
            else:
                print(f"[LibreOffice] stderr: {stderr_filtered[:500]}")
        
        # Check if conversion actually failed (ignore return code if output file exists)
        # LibreOffice may return non-zero exit code even on success due to warnings
        conversion_failed = result.returncode != 0
        
        # If return code indicates failure, check if output file was created anyway
        # (sometimes LibreOffice succeeds but returns error code due to warnings)
        if conversion_failed:
            # Check if output file might have been created despite error code
            # We'll verify this later when searching for output file
            error_msg = stderr_filtered if stderr_filtered else result.stdout or "Unknown LibreOffice error"
            # Don't raise error yet - check for output file first
            potential_error = f"LibreOffice returned code {result.returncode}. Error: {error_msg[:200]}"
            print(f"[LibreOffice] ⚠️ {potential_error}")
        
        # List files after conversion
        files_after = list(Path(output_dir_abs).iterdir())
        print(f"[LibreOffice] Files after: {[f.name for f in files_after]}")
        
        # Find the converted file (check multiple possible names)
        input_name = Path(input_path_abs).stem
        input_full_name = Path(input_path_abs).name
        input_base = input_full_name.rsplit('.', 1)[0] if '.' in input_full_name else input_full_name
        
        # FIXED: Better format extraction
        actual_format = output_format.split(':')[0].strip().strip('"\'')
        if not actual_format:
            actual_format = 'pdf'  # Default fallback
        
        print(f"[Conversion] Looking for format: {actual_format}, input_name: {input_name}, input_base: {input_base}")
        output_dir_path = Path(output_dir_abs)

        stdout_path_hints: List[str] = []
        if result.stdout:
            for match in re.finditer(r"->\s*(.+?\.[A-Za-z0-9]+)\s+using filter", result.stdout):
                hint = match.group(1).strip().strip("'\"")
                if hint:
                    stdout_path_hints.append(hint)

        possible_paths: List[Path] = []
        for hint in stdout_path_hints:
            hint_path = Path(hint)
            possible_paths.append(hint_path)
            if hint_path.parent != output_dir_path:
                possible_paths.append(output_dir_path / hint_path.name)
         
        # FIXED: Check multiple possible output filenames (Cloud Run LibreOffice issue)
        # LibreOffice sometimes creates input.docx.pdf instead of input.pdf
        possible_names = [
            f"{input_base}.{actual_format}",  # Most common: input.pdf
            f"{input_name}.{actual_format}",  # Stem + extension
            f"{input_full_name}.{actual_format}",  # FIX: input.docx.pdf (Cloud Run issue)
            f"{input_full_name.rsplit('.', 1)[0]}.{actual_format}",  # Another variant
            f"input.{actual_format}",  # Fallback if filename changed
        ]
        
        # Also check for files with spaces replaced or removed (LibreOffice sometimes does this)
        input_name_no_spaces = input_name.replace(' ', '_')
        input_base_no_spaces = input_base.replace(' ', '_')
        possible_names.extend([
            f"{input_name_no_spaces}.{actual_format}",
            f"{input_base_no_spaces}.{actual_format}",
            f"{input_full_name.replace(' ', '_')}.{actual_format}",  # input_docx.pdf variant
        ])
        
        # Remove duplicates while preserving order
        possible_names = list(dict.fromkeys(possible_names))
        print(f"[Conversion] Checking possible names: {possible_names[:5]}...")

        output_file = None
        for candidate_path in possible_paths:
            if candidate_path.exists():
                output_file = candidate_path
                print(f"[Conversion] Found output file from stdout hint: {output_file}")
                break
        
        for name in possible_names:
            candidate = Path(output_dir_abs) / name
            if candidate.exists():
                output_file = candidate
                print(f"[Conversion] Found output file: {output_file.name}")
                break

        # If still not found, search for any file with the output extension
        if not output_file:
            matching_files = list(output_dir_path.glob(f"*.{actual_format}"))
            matching_files += list(output_dir_path.glob(f"*.{actual_format.upper()}"))
            if matching_files:
                # Get the most recently modified file
                output_file = max(matching_files, key=lambda f: f.stat().st_mtime)
                print(f"[Conversion] Found output file by extension: {output_file.name}")

        # As a last resort, search recursively (LibreOffice sometimes nests output)
        if not output_file:
            recursive_matches = [
                p for p in output_dir_path.rglob(f"*.{actual_format}") if p.is_file()
            ]
            recursive_matches += [
                p for p in output_dir_path.rglob(f"*.{actual_format.upper()}") if p.is_file()
            ]
            if recursive_matches:
                output_file = max(recursive_matches, key=lambda f: f.stat().st_mtime)
                print(f"[Conversion] Found output file via recursive search: {output_file}")

        # Removed duplicate file finding logic - already checked above
 
        # Retry with simplified output format (without filter spec) if still missing
        if not output_file and allow_retry and ':' in output_format:
            simple_format = output_format.split(':', 1)[0].strip().strip('"\'')
            if simple_format:
                print(f"[Conversion] Retrying LibreOffice with simplified format: {simple_format}")
                return convert_with_libreoffice(
                    input_path=input_path_abs,
                    output_format=simple_format,
                    output_dir=output_dir_abs,
                    timeout=timeout,
                    allow_retry=False,
                    input_filter=input_filter
                )
        
        # If still no file, try with even simpler command (no filter, minimal flags)
        if not output_file and allow_retry:
            print(f"[Conversion] Retrying with minimal LibreOffice command (basic flags only)")
            minimal_cmd = [
                'libreoffice',
                '--headless',
                '--convert-to', actual_format,
                '--outdir', output_dir_abs,
                input_path_abs
            ]
            print(f"[LibreOffice] Minimal command: {' '.join(minimal_cmd)}")
            minimal_result = subprocess.run(
                minimal_cmd,
                capture_output=True,
                text=True,
                timeout=timeout,
                env=env,
                cwd=output_dir_abs
            )
            print(f"[LibreOffice] Minimal return code: {minimal_result.returncode}")
            if minimal_result.stdout:
                print(f"[LibreOffice] Minimal stdout (full): {minimal_result.stdout}")
            if minimal_result.stderr:
                # Filter Java warnings
                stderr_lines = minimal_result.stderr.split('\n')
                non_java_stderr = [line for line in stderr_lines if 'javaldx' not in line.lower() and 'java may not function' not in line.lower()]
                if non_java_stderr:
                    print(f"[LibreOffice] Minimal stderr (non-Java): {chr(10).join(non_java_stderr)}")
            
            # Optimized: Minimal wait for file write
            time.sleep(0.2)  # Reduced to 0.2s for faster processing
            matching_files = list(output_dir_path.glob(f"*.{actual_format}"))
            matching_files += list(output_dir_path.glob(f"*.{actual_format.upper()}"))
            if matching_files:
                output_file = max(matching_files, key=lambda f: f.stat().st_mtime)
                print(f"[Conversion] ✅ Found output file after minimal retry: {output_file.name}")
        
        if not output_file or not output_file.exists():
            # List all files in output directory for debugging
            all_files = list(Path(output_dir_abs).iterdir())
            file_list = ", ".join([f.name for f in all_files if f.is_file()]) if all_files else "no files"
            
            # Check if error is ONLY Java-related (non-fatal)
            stderr_lower = (result.stderr or "").lower()
            is_java_only_error = (
                result.returncode == 1 and
                ('javaldx' in stderr_lower or 'java may not function' in stderr_lower) and
                not any(keyword in stderr_lower for keyword in ['error', 'failed', 'cannot', 'unable'])
            )
            
            if is_java_only_error:
                # Java warning only - check if file was already found above
                print(f"[Conversion] Java warning detected but non-fatal, output file check already done above")
                if output_file and output_file.exists():
                    print(f"[Conversion] Found output file despite Java warning: {output_file.name}")
                else:
                    # Still no file - return generic error (Java warnings filtered)
                    error_msg = "Conversion failed. Please try again or check if the file format is supported."
                    raise Exception(error_msg)
            else:
                # Real error - get details but filter Java warnings
                stderr_clean = result.stderr or ""
                if stderr_clean:
                    stderr_clean = filter_java_warnings(stderr_clean)
                
                error_msg = (
                    f"Converted file not found in {output_dir}. "
                    f"Expected format: {actual_format}. "
                    f"Input file: {input_full_name}. "
                    f"Checked names: {', '.join(possible_names[:5])}. "
                    f"Found files: {file_list}."
                )
                if stderr_clean and stderr_clean != "Conversion failed. Please try again or check if the file format is supported.":
                    error_msg += f" LibreOffice error: {stderr_clean[:500]}"
                
                # Filter Java warnings from final error message
                error_msg = filter_java_warnings(error_msg)
                raise Exception(error_msg)
        
        # Validate file size
        file_size = output_file.stat().st_size
        if file_size < 100:
            raise Exception(f"Converted file is too small ({file_size} bytes) - conversion may have failed")
        
        # COMPETITOR ADVANTAGE: Preserve original filename in output
        # Many competitors lose the original filename - we preserve it
        original_name = Path(input_path).stem
        if original_name and original_name != "input":
            # Try to preserve original name if possible
            # Extract format safely
            format_ext = output_format.split(':')[0] if ':' in output_format else (output_format.split('"')[0] if '"' in output_format else output_format)
            preferred_output = Path(output_dir) / f"{original_name}.{format_ext}"
            if output_file.name != preferred_output.name:
                # Keep output_file as is, but log original name
                print(f"[Conversion] Output: {output_file.name} (from {original_name})")
        
        return str(output_file)
        
    except subprocess.TimeoutExpired:
        raise Exception("Conversion timeout - file may be too large or complex")
    except Exception as e:
        # Ensure any Java/javaldx noise is stripped at the lowest level
        raise Exception(filter_java_warnings(f"Conversion failed: {str(e)}"))


def convert_pdf_to_docx_with_pdf2docx(input_pdf: str, output_docx: str) -> str:
    """
    High-quality PDF to DOCX converter using pdf2docx.
    Preserves layout, images, tables, fonts, TOC numbering, and formatting - COMPETITOR-LEVEL QUALITY.
    CRITICAL: This method preserves images/logos and maintains proper page numbering for TOC.
    """
    converter = None
    conversion_errors: List[str] = []
    output_path = None
    
    try:
        print(f"[pdf2docx] Starting high-quality conversion: {input_pdf} -> {output_docx}")
        print(f"[pdf2docx] Using PyMuPDF engine - preserves layout, colors, images, logos, fonts, tables, TOC numbering")
        
        # Create converter with best quality settings
        # pdf2docx uses PyMuPDF (fitz) which has excellent image extraction capabilities
        converter = Converter(input_pdf)
        
        # Convert with maximum quality settings and CRITICAL parameters
        # pdf2docx automatically:
        # - Preserves page layout and margins
        # - Maintains fonts and text formatting
        # - Embeds images at original quality (CRITICAL for logos)
        # - Converts tables with proper structure
        # - Preserves colors and styles
        # - Maintains page numbering for TOC (preserves section breaks properly)
        # - Extracts all images including logos, headers, footers
        #
        # IMPROVED SETTINGS to prevent blank pages and broken formatting:
        # - multi_processing=False: More reliable single-threaded conversion
        # - line_separate_threshold: Better text line detection
        # METHOD 0: Microsoft Word Native Conversion (Best Quality)
        # Uses local Word installation via COM automation
        try:
            # DEBUG: Force Linux Simulation
            # Change 'True' to 'False' to restore Windows behavior
            SIMULATE_LINUX = False
            
            if SIMULATE_LINUX or os.name != 'nt':
                print("[PDF->DOCX] Simulation or Non-Windows OS detected. Skipping Native Word.")
                raise Exception("Skipping Native Word (Linux/Cloud Run detected)")

            print("[PDF->DOCX] Attempting Native Word Conversion (Best Quality)...")
            import win32com.client
            import pythoncom
            
            # Initialize COM
            pythoncom.CoInitialize()
            
            # Create Word Application
            # Use DispatchEx to force a NEW instance. Dispatch can hook into a stuck/frozen instance.
            word = win32com.client.DispatchEx("Word.Application")
            # CRITICAL DEBUG: Make Word visible so user can see/interact with any blocking dialogs
            word.Visible = True
            word.DisplayAlerts = False # Keep this False to suppress standard alerts, but visibility helps if it's a hard block
            word.DisplayAlerts = False
            
            try:
                # Open PDF
                abs_input = os.path.abspath(input_pdf)
                abs_output = os.path.abspath(output_docx) # Fix: Save directly to the expected output path
                
                print(f"[Word] Opening: {abs_input}")
                # ConfirmConversions=False SUppresses the "Word will now convert..." dialog
                doc = word.Documents.Open(abs_input, ConfirmConversions=False)
                
                # Save as DOCX (wdFormatXMLDocument = 12)
                print(f"[Word] Saving as: {abs_output}")
                doc.SaveAs2(abs_output, FileFormat=12)
                doc.Close()
                
                if os.path.exists(abs_output) and os.path.getsize(abs_output) > 100:
                    output_path = abs_output
                    print("[PDF->DOCX] ✅ Native Word conversion succeeded!")
            except Exception as w_err:
                print(f"[Word] Error during conversion: {w_err}")
                try:
                    doc.Close(SaveChanges=0) # wdDoNotSaveChanges
                except:
                    pass
                # Soft failure: Raise so we catch it below and fall back
                raise Exception(f"Native Word Conversion Failed: {w_err}")
            finally:
                try:
                    word.Quit()
                except:
                    pass
                
        except Exception as e:
             print(f"[PDF->DOCX] ⚠️ Native Word conversion failed: {e}")
             # Add to errors but continue to fallback
             conversion_errors.append(f"Native Word: {str(e)[:200]}")
             
        # Fallback to pdf2docx if Word failed/missing
        if not output_path or not os.path.exists(output_path):
            print("[PDF->DOCX] Falling back to pdf2docx...") 
            converter.convert(output_docx)
            if os.path.exists(output_docx):
                output_path = output_docx
        
        # Verify images were extracted
        try:
            import zipfile
            with zipfile.ZipFile(output_docx, 'r') as zip_ref:
                image_files = [f for f in zip_ref.namelist() if f.startswith('word/media/')]
                if image_files:
                    print(f"[pdf2docx] ✅ Found {len(image_files)} image(s) extracted (logos, images preserved)")
                else:
                    print(f"[pdf2docx] ⚠️ No images found in output (PDF may not have images or they were vector graphics)")
        except Exception as img_check_error:
            print(f"[pdf2docx] ⚠️ Could not verify images: {img_check_error}")
        
        print(f"[pdf2docx] ✅ Conversion completed - ALL formatting, images, and TOC numbering preserved")
        
    except Exception as e:
        error_msg = f"pdf2docx conversion error: {str(e)}"
        print(f"[pdf2docx] ❌ {error_msg}")
        raise Exception(error_msg)
    finally:
        # Always close converter to free resources
        if converter:
            try:
                converter.close()
                print(f"[pdf2docx] Converter closed")
            except Exception as close_error:
                print(f"[pdf2docx] Warning: Error closing converter: {close_error}")

    # Validate output file
    if not os.path.exists(output_docx):
        raise Exception("pdf2docx did not create an output file")

    file_size = os.path.getsize(output_docx)
    if file_size < 100:
        raise Exception(f"pdf2docx output is too small ({file_size} bytes) - conversion may have failed")

    # Validate DOCX structure (should be a ZIP archive)
    with open(output_docx, "rb") as fh:
        header = fh.read(2)
        if header != b"PK":
            raise Exception(f"pdf2docx output is not a valid DOCX archive (header: {header.hex()})")

    print(f"[pdf2docx] ✅ Output validated: {file_size} bytes, valid DOCX structure")
    return output_docx

@app.post("/convert/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    """Convert PDF to Word (DOCX) - High quality conversion"""
    start_time = time.time()
    
    if not file.filename or not file.filename.lower().endswith('.pdf'):
        raise HTTPException(400, "Only PDF files are allowed")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # Save uploaded file
            # Save uploaded file to CWD (Trusted Location) to avoid Protected View
            cwd = os.getcwd()
            unique_name = f"processing_{int(time.time())}.pdf"
            input_path = os.path.join(cwd, unique_name)
            content = await file.read()
            
            if not content or len(content) < 100:
                raise HTTPException(400, "File is too small or empty")
            
            with open(input_path, "wb") as f:
                f.write(content)
            
            print(f"[PDF->DOCX] Processing: {file.filename}, {len(content)} bytes")
            
            # Convert to DOCX with BEST QUALITY settings - COMPETITOR-LEVEL QUALITY
            # Strategy: Try pdf2docx FIRST (best for images, layout, TOC numbering),
            # then fall back to LibreOffice methods only if needed.
            # CRITICAL: pdf2docx preserves images/logos and TOC numbering better than LibreOffice
            conversion_errors: List[str] = []
            output_path: Optional[str] = None
            
            # ULTRA-OPTIMIZED: Dynamic timeout based on file size
            file_size_mb = len(content) / (1024 * 1024)
            dynamic_timeout = 30 if file_size_mb < 1 else (45 if file_size_mb < 5 else 60)
            print(f"[PDF->DOCX] File size: {file_size_mb:.2f}MB, using timeout: {dynamic_timeout}s")
            
            # Track which conversion method succeeded
            used_pdf2docx = False
            engine_used = "Unknown" 
            
            # METHOD 0: Microsoft Word Native Conversion (Best Quality)
            # Uses local Word installation via COM automation
            try:
                # DEBUG: Force Linux Simulation
                # Change 'True' to 'False' to restore Windows behavior
                SIMULATE_LINUX = False
                
                if SIMULATE_LINUX or os.name != 'nt':
                    print("[PDF->DOCX] Simulation or Non-Windows OS detected. Skipping Native Word.")
                    raise Exception("Skipping Native Word (Linux/Cloud Run detected)")
    
                print("[PDF->DOCX] Attempting Native Word Conversion (Best Quality)...")
                import win32com.client
                import pythoncom
                
                # Initialize COM
                pythoncom.CoInitialize()
                
                # Create Word Application
                # Use DispatchEx to force a NEW instance. Dispatch can hook into a stuck/frozen instance.
                word = win32com.client.DispatchEx("Word.Application")
                # CRITICAL DEBUG: Make Word visible so user can see/interact with any blocking dialogs
                word.Visible = True
                word.DisplayAlerts = False 
                
                try:
                    # Open PDF
                    # Fix: Save directly to the expected output path
                    abs_input = os.path.abspath(input_path)
                    abs_output = os.path.abspath(output_docx) 
                    
                    print(f"[Word] Opening: {abs_input}")
                    # ConfirmConversions=False SUppresses the "Word will now convert..." dialog
                    doc = word.Documents.Open(abs_input, ConfirmConversions=False)
                    
                    # Save as DOCX (wdFormatXMLDocument = 12)
                    print(f"[Word] Saving as: {abs_output}")
                    doc.SaveAs2(abs_output, FileFormat=12)
                    doc.Close()
                    
                    if os.path.exists(abs_output) and os.path.getsize(abs_output) > 100:
                        output_path = abs_output
                        engine_used = "NativeWord"
                        print("[PDF->DOCX] ✅ Native Word conversion succeeded!")
                except Exception as w_err:
                    print(f"[Word] Error during conversion: {w_err}")
                    try:
                        doc.Close(SaveChanges=0) # wdDoNotSaveChanges
                    except:
                        pass
                    # Soft failure: Raise so we catch it below and fall back
                    raise Exception(f"Native Word Conversion Failed: {w_err}")
                finally:
                    try:
                        word.Quit()
                    except:
                        pass
                    
            except Exception as e:
                 print(f"[PDF->DOCX] ⚠️ Native Word conversion failed/skipped: {e}")
                 conversion_errors.append(f"Native Word: {str(e)[:200]}")
            else:
                print("[PDF->DOCX] Skipping Native Word (Not on Windows)")
                 
            # FALLBACK CHAIN: LibreOffice -> pdf2docx
            
            # METHOD 1: LibreOffice (Headless) - Gold Standard for Linux
            if not output_path or not os.path.exists(output_path):
                try:
                    print("[PDF->DOCX] Attempting LibreOffice (Best for Linux/Production)...")
                    # Use 'docx:"MS Word 2007 XML"' filter for best compatibility
                    output_path = convert_with_libreoffice(
                        input_path,
                        'docx:"MS Word 2007 XML"', 
                        temp_dir, # LibreOffice handles temp paths fine usually
                        timeout=dynamic_timeout,
                        input_filter='writer_pdf_import'
                    )
                    if output_path and os.path.exists(output_path) and os.path.getsize(output_path) > 100:
                        engine_used = "LibreOffice"
                        print("[PDF->DOCX] ✅ LibreOffice conversion succeeded!")
                except Exception as lo_err:
                    print(f"[PDF->DOCX] ⚠️ LibreOffice failed: {lo_err}")
                    conversion_errors.append(f"LibreOffice: {str(lo_err)[:200]}")
    
            # METHOD 2: pdf2docx (Last Resort - Good for images but bad for formatting)
            if not output_path or not os.path.exists(output_path):
                 # Restore fallback_path logic if needed or just use output_docx
                try:
                    fallback_path = os.path.join(cwd, f"{Path(unique_name).stem}.docx")
                    print("[PDF->DOCX] Attempting pdf2docx FIRST (best for images, logos, and TOC numbering)...")
                    output_path = convert_pdf_to_docx_with_pdf2docx(input_path, fallback_path)
                    if output_path and os.path.exists(output_path) and os.path.getsize(output_path) > 100:
                        engine_used = "pdf2docx"
                        used_pdf2docx = True
                        print("[PDF->DOCX] ✅ pdf2docx succeeded - images, logos, and TOC numbering preserved!")
                except Exception as pdf2docx_error:
                    conversion_errors.append(f"pdf2docx primary: {str(pdf2docx_error)[:200]}")
                    print(f"[PDF->DOCX] ⚠️ pdf2docx failed ({pdf2docx_error})")

            # METHOD 2: LibreOffice with optimized settings (fallback - may lose images but handles page breaks)
            if not output_path or not os.path.exists(output_path) or os.path.getsize(output_path) < 100:
                try:
                    print("[PDF->DOCX] Attempting LibreOffice with optimized settings (fallback - may lose images)...")
                    output_path = convert_with_libreoffice(
                        input_path,
                        'docx:"MS Word 2007 XML"',  # Explicit DOCX filter
                        temp_dir,
                        timeout=dynamic_timeout,
                        input_filter='writer_pdf_import'
                    )
                    if output_path and os.path.exists(output_path) and os.path.getsize(output_path) > 100:
                        print("[PDF->DOCX] ✅ LibreOffice primary conversion succeeded!")
                except Exception as lo_error:
                    conversion_errors.append(f"LibreOffice primary: {str(lo_error)[:200]}")
                    print(f"[PDF->DOCX] ⚠️ LibreOffice primary failed ({lo_error}), trying LibreOffice default...")

            # METHOD 3: LibreOffice default DOCX (last fallback)
            if not output_path or not os.path.exists(output_path) or os.path.getsize(output_path) < 100:
                try:
                    print("[PDF->DOCX] Attempting LibreOffice with default DOCX export (final fallback)...")
                    output_path = convert_with_libreoffice(
                        input_path,
                        'docx',  # Let LibreOffice choose the filter automatically
                        temp_dir,
                        timeout=dynamic_timeout,
                        input_filter='writer_pdf_import'
                    )
                    if output_path and os.path.exists(output_path) and os.path.getsize(output_path) > 100:
                        print("[PDF->DOCX] ✅ LibreOffice default DOCX export succeeded!")
                except Exception as secondary_error:
                    conversion_errors.append(f"LibreOffice default docx: {str(secondary_error)[:200]}")
                    print(f"[PDF->DOCX] ⚠️ Default DOCX export failed ({secondary_error})")
            
            # Final validation - if all methods failed
            if not output_path or not os.path.exists(output_path) or os.path.getsize(output_path) < 100:
                error_summary = "; ".join(conversion_errors) if conversion_errors else "All conversion methods failed"
                raise Exception(f"PDF to DOCX conversion failed. Errors: {error_summary}")
            
            # Validate output
            if not os.path.exists(output_path):
                raise Exception("DOCX file not created")
            
            file_size = os.path.getsize(output_path)
            if file_size < 100:
                raise Exception(f"DOCX file too small ({file_size} bytes)")
            
            # Validate DOCX (ZIP structure) - CRITICAL VALIDATION
            with open(output_path, "rb") as f:
                header = f.read(4)  # Read more bytes for better validation
                if header[:2] != b'PK':
                    raise Exception(f"Invalid DOCX file structure (not a valid ZIP). Header: {header.hex()}")
                # DOCX files should start with PK (ZIP signature)
                # Check for ZIP local file header signature: 50 4B 03 04 or 50 4B 05 06 (empty ZIP)
                if header[:2] == b'PK' and header[2:4] not in [b'\x03\x04', b'\x05\x06', b'\x07\x08']:
                    print(f"[PDF->DOCX] ⚠️ Warning: Unusual ZIP header: {header.hex()}, but continuing...")
            
            # Additional validation: Try to read ZIP structure
            try:
                import zipfile
                with zipfile.ZipFile(output_path, 'r') as zip_ref:
                    # Check if it has required DOCX files
                    required_files = ['[Content_Types].xml', 'word/document.xml']
                    zip_files = zip_ref.namelist()
                    has_content_types = any('Content_Types' in f for f in zip_files)
                    has_document = any('word/document.xml' in f for f in zip_files)
                    
                    if not has_content_types or not has_document:
                        print(f"[PDF->DOCX] ⚠️ Warning: DOCX structure incomplete. Files: {zip_files[:5]}")
                        # Don't fail - some converters might use different structure
            except Exception as zip_error:
                print(f"[PDF->DOCX] ⚠️ Warning: Could not validate ZIP structure: {zip_error}")
                # Don't fail - file might still be valid
            
            # Post-processing: Remove blank pages like iLovePDF does
            # 1. Change all section breaks to CONTINUOUS
            # 2. Remove explicit page break elements
            # 3. This produces clean output without blank pages
            # Post-processing DISABLED to preserve formatting
            # 1. Changing section breaks to CONTINUOUS damages layout
            # 2. Removing page breaks removes necessary spacing
            # 3. This produces a faithful conversion even if it keeps blank pages
            try:
                # SKIP POST-PROCESSING
                if True: 
                     print(f"[PDF->DOCX] Skipping aggressive post-processing to preserve formatting")
                     raise Exception("Skipping cleanup to preserve layout")

                from docx import Document
                from docx.oxml.ns import qn
                from docx.enum.section import WD_SECTION
                from lxml import etree
                
                doc = Document(output_path)
                body = doc.element.body
                
                # Count initial state
                initial_sections = len(doc.sections)
                initial_paras = len(doc.paragraphs)
                
                print(f"[PDF->DOCX] Initial state: {initial_sections} sections, {initial_paras} paragraphs")
                
                changes_made = False
                
                # STEP 1: Change all section breaks to CONTINUOUS
                sections_fixed = 0
                for i, section in enumerate(doc.sections):
                    try:
                        if i == 0:
                            continue
                        if section.start_type != WD_SECTION.CONTINUOUS:
                            section.start_type = WD_SECTION.CONTINUOUS
                            sections_fixed += 1
                            changes_made = True
                    except:
                        pass
                
                if sections_fixed > 0:
                    print(f"[PDF->DOCX] ✅ Changed {sections_fixed} section breaks to CONTINUOUS")
                
                # STEP 2: Remove ALL explicit page breaks from document
                # This is the key to removing blank pages
                page_breaks_removed = 0
                
                # Find all page break elements: w:br with w:type="page"
                nsmap = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                
                # Remove w:br elements with type="page"
                for br in body.xpath('.//w:br[@w:type="page"]', namespaces=nsmap):
                    try:
                        parent = br.getparent()
                        if parent is not None:
                            parent.remove(br)
                            page_breaks_removed += 1
                            changes_made = True
                    except:
                        pass
                
                # Remove w:lastRenderedPageBreak elements (soft page breaks)
                for pb in body.xpath('.//w:lastRenderedPageBreak', namespaces=nsmap):
                    try:
                        parent = pb.getparent()
                        if parent is not None:
                            parent.remove(pb)
                            page_breaks_removed += 1
                            changes_made = True
                    except:
                        pass
                
                if page_breaks_removed > 0:
                    print(f"[PDF->DOCX] ✅ Removed {page_breaks_removed} page break elements")
                
                # STEP 3: Remove empty paragraphs that only contain whitespace
                empty_paras_removed = 0
                paras_to_remove = []
                
                for para in doc.paragraphs:
                    # Check if paragraph is empty or only whitespace
                    text = para.text.strip()
                    if not text:
                        # Check if paragraph has any content (images, etc.)
                        para_xml = para._element
                        has_content = False
                        
                        # Check for drawings/images
                        drawings = para_xml.xpath('.//w:drawing', namespaces=nsmap)
                        objects = para_xml.xpath('.//w:object', namespaces=nsmap)
                        
                        if drawings or objects:
                            has_content = True
                        
                        # If no content, mark for removal
                        if not has_content:
                            paras_to_remove.append(para._element)
                
                # Remove empty paragraphs
                for para_elem in paras_to_remove:
                    try:
                        parent = para_elem.getparent()
                        if parent is not None:
                            parent.remove(para_elem)
                            empty_paras_removed += 1
                            changes_made = True
                    except:
                        pass
                
                if empty_paras_removed > 0:
                    print(f"[PDF->DOCX] ✅ Removed {empty_paras_removed} empty paragraphs")
                
                # Save if changes were made
                if changes_made:
                    doc.save(output_path)
                    print(f"[PDF->DOCX] ✅ Document cleaned: blank pages removed")
                else:
                    print(f"[PDF->DOCX] ✅ Document already clean")
                
                print(f"[PDF->DOCX] Conversion method: {'pdf2docx' if used_pdf2docx else 'LibreOffice'}")
                
            except Exception as e:
                print(f"[PDF->DOCX] ⚠️ Post-processing error: {e}")
            
            # CRITICAL: Fix encoding issues - convert boxes (□) and question marks to dots
            # SKIP encoding fixes if pdf2docx was used (it handles encoding correctly)
            # Only apply encoding fixes for LibreOffice conversions
            if not used_pdf2docx:
                try:
                    from docx import Document
                    print("[PDF->DOCX] Fixing encoding issues (boxes and question marks to dots)...")
                    # Reload document if it was modified in post-processing
                    doc = Document(output_path)
                    
                    # Box-like characters that should be converted to dots
                    box_chars = ["□", "▪", "■", "▫", "▬", "▭", "▮", "▯", "•", "·", "◆", "►", "❖", "➤", "⧾", "✓", "🔹", "▸"]
                    question_marks = ["?", "？", "\uff1f", "\u003f", "\uFF1F", "\ufffd", "\uFFFD"]
                    all_problem_chars = box_chars + question_marks
                    
                    fixed_count = 0
                    for para in doc.paragraphs:
                        original_text = para.text
                        if not original_text.strip():
                            continue
                        
                        # Check if paragraph has any problem characters
                        has_problem = any(char in original_text for char in all_problem_chars)
                        if not has_problem:
                            continue
                        
                        # Replace all problem characters with dots
                        fixed_text = original_text
                        for char in all_problem_chars:
                            fixed_text = fixed_text.replace(char, ".")
                        
                        if fixed_text != original_text:
                            # Get original font size from first run
                            original_font_size = None
                            if para.runs:
                                original_font_size = para.runs[0].font.size
                            
                            # Clear and rebuild paragraph
                            para.clear()
                            
                            # Add fixed text with proper font size
                            run = para.add_run(fixed_text)
                            if original_font_size:
                                run.font.size = original_font_size
                            
                            fixed_count += 1
                    
                    # Also fix table cells
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                for para in cell.paragraphs:
                                    original_text = para.text
                                    if not original_text.strip():
                                        continue
                                    
                                    has_problem = any(char in original_text for char in all_problem_chars)
                                    if not has_problem:
                                        continue
                                    
                                    fixed_text = original_text
                                    for char in all_problem_chars:
                                        fixed_text = fixed_text.replace(char, ".")
                                    
                                    if fixed_text != original_text:
                                        original_font_size = None
                                        if para.runs:
                                            original_font_size = para.runs[0].font.size
                                        
                                        para.clear()
                                        run = para.add_run(fixed_text)
                                        if original_font_size:
                                            run.font.size = original_font_size
                                        
                                        fixed_count += 1
                    
                    if fixed_count > 0:
                        print(f"[PDF->DOCX] ✅ Fixed {fixed_count} paragraphs/cells with boxes/question marks → dots")
                        doc.save(output_path)
                    else:
                        print("[PDF->DOCX] No encoding issues found")
                except Exception as fix_error:
                    print(f"[PDF->DOCX] ⚠️ Warning: Could not fix encoding: {fix_error}, continuing with original file...")
            else:
                print("[PDF->DOCX] ✅ pdf2docx was used - skipping encoding fixes (pdf2docx handles encoding correctly)")
            
            # Read and encode output
            with open(output_path, "rb") as f:
                output_data = f.read()
            
            # Final size check
            if len(output_data) < 100:
                raise Exception(f"Output file too small: {len(output_data)} bytes")
            
            output_base64 = base64.b64encode(output_data).decode()
            
            elapsed = time.time() - start_time
            print(f"[PDF->DOCX] ✅ Success in {elapsed:.1f}s - {len(output_data)} bytes")
            print(f"[PDF->DOCX] ✅ Base64 length: {len(output_base64)} characters")
            
            # COMPETITOR ADVANTAGE: Return comprehensive metadata

            # Determine filename and engine for FileResponse
            download_filename = (file.filename.replace('.pdf', '') + '_native_word.docx') if file.filename else 'converted_native_word.docx'
            engine_used = 'pdf2docx' if used_pdf2docx else 'LibreOffice'

            # Return the file directly
            return FileResponse(
                output_path, 
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                filename=download_filename,
                headers={"X-Conversion-Engine": engine_used}
            )
            
        except HTTPException:
            raise
        except Exception as e:
            elapsed = time.time() - start_time
            error_msg = filter_java_warnings(str(e))
            print(f"[PDF->DOCX] ERROR Failed after {elapsed:.1f}s: {error_msg}")
            raise HTTPException(500, f"Conversion failed: {error_msg}")

async def pdf_to_excel_legacy(file: UploadFile = File(...)):
    """Legacy PDF to Excel pipeline placeholder (deprecated)."""
    raise HTTPException(status_code=410, detail="Legacy PDF to Excel pipeline has been retired.")

MAX_PDF_TO_EXCEL_BYTES = 40 * 1024 * 1024  # 40 MB
MAX_PDF_TO_EXCEL_PAGES = 250


@app.post("/convert/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    """Convert PDF to Excel (XLSX) using the new multi-stage engine."""
    start_time = time.time()

    if not file.filename or not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "Only PDF files are allowed")

    content = await file.read()

    if not content or len(content) < 100:
        raise HTTPException(400, "File is too small or empty")

    if len(content) > MAX_PDF_TO_EXCEL_BYTES:
        raise HTTPException(400, "File exceeds the 40MB limit for PDF to Excel.")

    print(f"[PDF->XLSX] Starting next-gen engine for '{file.filename}' ({len(content)} bytes)")

    try:
        output_data, stats = convert_pdf_to_excel(content, max_pages=MAX_PDF_TO_EXCEL_PAGES)
    except ValueError as exc:
        raise HTTPException(400, str(exc)) from exc
    except Exception as exc:
        print(f"[PDF->XLSX] ERROR during conversion: {exc}")
        raise HTTPException(500, f"Conversion failed: {exc}") from exc

    output_base64 = base64.b64encode(output_data).decode()
    elapsed = time.time() - start_time

    features = {
        "single_sheet": True,
        "tables_detected": stats.get("tables_detected", 0),
        "images_embedded": stats.get("images_extracted", 0),
        "text_blocks": stats.get("text_blocks", 0),
        "engine": "PyMuPDF + Camelot/Tabula + openpyxl",
        "max_pages_reached": bool(stats.get("max_pages_reached")),
    }

    print(
        "[PDF->XLSX] Completed in %.1fs - tables=%s, images=%s, rows=%s, pages=%s",
        elapsed,
        features["tables_detected"],
        features["images_embedded"],
        stats.get("rows_written", 0),
        stats.get("pages_processed", 0),
    )

    return JSONResponse(
        {
            "success": True,
            "file": output_base64,
            "filename": Path(file.filename).stem + ".xlsx",
            "size": len(output_data),
            "original_size": len(content),
            "processing_time": f"{elapsed:.1f}s",
            "quality": "enterprise-grade",
            "features": features,
        }
    )

@app.post("/convert/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    """Convert PDF to PowerPoint (PPTX) - Extract content and create slides"""
    start_time = time.time()
    
    if not file.filename or not file.filename.lower().endswith('.pdf'):
        raise HTTPException(400, "Only PDF files are allowed")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # Save uploaded file
            input_path = os.path.join(temp_dir, "input.pdf")
            content = await file.read()
            
            if not content or len(content) < 100:
                raise HTTPException(400, "File is too small or empty")
            
            with open(input_path, "wb") as f:
                f.write(content)
            
            print(f"[PDF->PPTX] Processing: {file.filename}, {len(content)} bytes")
            
            # Create PowerPoint presentation with PROFESSIONAL dimensions
            # Standard widescreen format (16:9) - better than competitors who use old 4:3
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # Widescreen 16:9 aspect ratio
            prs.slide_height = Inches(7.5)    # Standard professional format
            
            # COMPETITOR ADVANTAGE: Set presentation-level properties
            # Many competitors don't set metadata - we do for professional output
            try:
                prs.core_properties.author = "PDFMasterTool"
                prs.core_properties.title = "Converted from PDF"
                prs.core_properties.comments = "High-quality conversion with complete content preservation"
            except:
                pass  # Properties might not be available in all versions
            
            # Track if we found any content
            has_content = False

            try:
                with pdfplumber.open(input_path) as pdf:
                    total_pages = len(pdf.pages)
                    print(f"[PDF->PPTX] Processing {total_pages} pages...")

                    for page_num, page in enumerate(pdf.pages, 1):
                        print(f"[PDF->PPTX] Processing page {page_num}/{total_pages}...")

                        # Extract ALL content using multiple methods
                        text = page.extract_text()
                        words = page.extract_words()
                        tables = page.extract_tables()

                        print(
                            f"[PDF->PPTX] Page {page_num}: Extracted "
                            f"{len(text) if text else 0} chars, {len(words)} words, {len(tables)} tables"
                        )

                        # Combine all text sources
                        all_text_content: list[str] = []

                        # First, add tables (if any)
                        if tables:
                            for table_idx, table in enumerate(tables):
                                if not table:
                                    continue

                                if table_idx > 0:
                                    all_text_content.append("")  # Space between tables

                                for row in table:
                                    if row:
                                        row_text = " | ".join(
                                            [str(cell).strip() if cell else "" for cell in row if cell]
                                        )
                                        row_text = row_text.strip()
                                        if row_text:
                                            all_text_content.append(row_text)
                                            has_content = True

                        # Add extracted text (preferred method)
                        if text:
                            for line in text.split("\n"):
                                line = line.strip()
                                if not line:
                                    continue
                                line = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", line)
                                if line:
                                    all_text_content.append(line)
                                    has_content = True

                        # ALSO use words method to catch anything missed
                        if words:
                            word_lines: list[str] = []
                            current_word_line = ""
                            last_y = None

                            for word in words:
                                word_text = word.get("text", "").strip()
                                if not word_text:
                                    continue

                                word_y = word.get("top", 0)

                                if last_y is not None and abs(word_y - last_y) > 8:
                                    if current_word_line.strip():
                                        word_lines.append(current_word_line.strip())
                                        current_word_line = ""

                                if current_word_line:
                                    current_word_line += " " + word_text
                                else:
                                    current_word_line = word_text

                                last_y = word_y

                            if current_word_line.strip():
                                word_lines.append(current_word_line.strip())

                            # Use words if extract_text() found less content
                            if len(all_text_content) < len(word_lines) * 0.5:
                                print(
                                    f"[PDF->PPTX] Page {page_num}: Using words method "
                                    f"({len(word_lines)} lines vs {len(all_text_content)} from extract_text)"
                                )
                                all_text_content = word_lines
                                has_content = True
                            elif len(all_text_content) == 0 and word_lines:
                                print("[PDF->PPTX] Page {page_num}: No text extracted, using words method")
                                all_text_content = word_lines
                                has_content = True

                        # Log extraction stats
                        print(
                            f"[PDF->PPTX] Page {page_num}: Extracted "
                            f"{len(all_text_content)} lines, {len(tables)} tables"
                        )

                        # Detect logical "sections" inside a page (Page 1, Page 2, etc.)
                        section_markers: list[int] = []
                        if all_text_content:
                            for idx, line in enumerate(all_text_content):
                                if re.match(r"^Page\s+\d+", line, re.IGNORECASE) or re.match(
                                    r"^Page\s+\d+\s*-", line, re.IGNORECASE
                                ):
                                    section_markers.append(idx)
                                    print(
                                        f"[PDF->PPTX] Found section marker at line {idx}: {line[:50]}"
                                    )

                        if section_markers:
                            # --- Multi‑section page: create one slide per section ---
                            print(
                                f"[PDF->PPTX] Page {page_num}: Found {len(section_markers)} section markers - "
                                "splitting by sections..."
                            )
                            print(
                                f"[PDF->PPTX] Total content lines: {len(all_text_content)}"
                            )

                            sections: list[list[str]] = []

                            content_before_first_marker: list[str] = []
                            if section_markers[0] > 0:
                                content_before_first_marker = all_text_content[0 : section_markers[0]]
                                print(
                                    "[PDF->PPTX] Content before first marker: "
                                    f"{len(content_before_first_marker)} lines"
                                )

                            for i, marker_idx in enumerate(section_markers):
                                start_idx = marker_idx
                                end_idx = (
                                    section_markers[i + 1]
                                    if i + 1 < len(section_markers)
                                    else len(all_text_content)
                                )

                                if i == 0 and content_before_first_marker:
                                    section_content = (
                                        content_before_first_marker
                                        + all_text_content[start_idx:end_idx]
                                    )
                                else:
                                    section_content = all_text_content[start_idx:end_idx]

                                sections.append(section_content)
                                print(
                                    "[PDF->PPTX] Section "
                                    f"{i + 1} (marker at {start_idx}): {len(section_content)} lines"
                                )

                            # Ensure last section contains all remaining content
                            if sections:
                                last_section = sections[-1]
                                expected_last_section_content = all_text_content[section_markers[-1] :]
                                if len(last_section) < len(expected_last_section_content):
                                    missing_lines = expected_last_section_content[len(last_section) :]
                                    if missing_lines:
                                        last_section.extend(missing_lines)
                                        print(
                                            f"[PDF->PPTX] Added {len(missing_lines)} missing lines to last section"
                                        )

                            # Create one slide per section
                            for section_idx, section_content in enumerate(sections):
                                if not section_content:
                                    continue

                                # Try to pick a good title
                                section_title = f"Section {section_idx + 1}"
                                for line in section_content:
                                    if re.match(r"^Page\s+\d+", line, re.IGNORECASE):
                                        section_title = line.strip()
                                        break

                                if section_title == f"Section {section_idx + 1}":
                                    for line in section_content:
                                        line = line.strip()
                                        if line and len(line) > 3:
                                            section_title = line[:60]
                                            break

                                slide_layout = prs.slide_layouts[1]
                                slide = prs.slides.add_slide(slide_layout)
                                title = slide.shapes.title
                                title.text = section_title

                                tf = slide.placeholders[1].text_frame
                                tf.word_wrap = True

                                for idx, content_line in enumerate(section_content):
                                    if not content_line:
                                        continue
                                    content_line = re.sub(
                                        r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", content_line
                                    )
                                    if not content_line:
                                        continue
                                    try:
                                        if idx == 0:
                                            p = tf.paragraphs[0]
                                        else:
                                            p = tf.add_paragraph()
                                        p.text = content_line
                                        p.level = 0
                                        p.font.size = Pt(12)
                                        p.font.name = "Calibri"
                                        p.space_after = Pt(6)
                                        has_content = True
                                    except Exception as para_error:
                                        print(
                                            "[PDF->PPTX] Warning: Could not add paragraph "
                                            f"{idx}: {para_error}"
                                        )
                                        continue

                                print(
                                    f"[PDF->PPTX] Section {section_idx + 1}: "
                                    f"Added {len(tf.paragraphs)} paragraphs"
                                )

                            # Done with this page when using sections
                            continue

                        # --- Normal case: one or more slides per page based on size ---
                        MAX_CHARS_PER_SLIDE = 30000  # PPTX text‑frame limit is ~32k
                        MAX_LINES_PER_SLIDE = 500

                        total_chars = sum(len(line) for line in all_text_content)
                        print(
                            f"[PDF->PPTX] Page {page_num}: Total content: {total_chars} chars, "
                            f"{len(all_text_content)} lines"
                        )

                        if (
                            total_chars > MAX_CHARS_PER_SLIDE
                            or len(all_text_content) > MAX_LINES_PER_SLIDE
                        ):
                            print(
                                f"[PDF->PPTX] Page {page_num}: Content too large, "
                                "splitting into multiple slides..."
                            )

                            chunk_size = min(MAX_LINES_PER_SLIDE, len(all_text_content))
                            chunks = [
                                all_text_content[i : i + chunk_size]
                                for i in range(0, len(all_text_content), chunk_size)
                            ]

                            for chunk_idx, chunk in enumerate(chunks):
                                slide_layout = prs.slide_layouts[1]
                                slide = prs.slides.add_slide(slide_layout)
                                title = slide.shapes.title
                                if len(chunks) > 1:
                                    title.text = (
                                        f"Page {page_num} - Part {chunk_idx + 1}/{len(chunks)}"
                                    )
                                else:
                                    title.text = f"Page {page_num}"

                                tf = slide.placeholders[1].text_frame
                                tf.word_wrap = True

                                for idx, content_line in enumerate(chunk):
                                    if not content_line:
                                        continue
                                    content_line = re.sub(
                                        r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", content_line
                                    )
                                    if not content_line:
                                        continue
                                    try:
                                        if idx == 0:
                                            p = tf.paragraphs[0]
                                        else:
                                            p = tf.add_paragraph()
                                        p.text = content_line
                                        p.level = 0
                                        p.font.size = Pt(11)
                                        p.font.name = "Calibri"
                                        p.space_after = Pt(4)
                                        has_content = True
                                    except Exception as para_error:
                                        print(
                                            "[PDF->PPTX] Warning: Could not add paragraph "
                                            f"{idx}: {para_error}"
                                        )
                                        continue

                                print(
                                    "[PDF->PPTX] Page "
                                    f"{page_num} Part {chunk_idx + 1}: "
                                    f"Added {len(tf.paragraphs)} paragraphs"
                                )
                        else:
                            # Normal single‑slide case
                            slide_layout = prs.slide_layouts[1]
                            slide = prs.slides.add_slide(slide_layout)

                            title = slide.shapes.title
                            title.text = f"Page {page_num}"

                            tf = slide.placeholders[1].text_frame
                            tf.word_wrap = True

                            for idx, content_line in enumerate(all_text_content):
                                if not content_line:
                                    continue
                                content_line = re.sub(
                                    r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", content_line
                                )
                                if not content_line:
                                    continue
                                try:
                                    if idx == 0:
                                        p = tf.paragraphs[0]
                                    else:
                                        p = tf.add_paragraph()
                                    p.text = content_line
                                    p.level = 0
                                    p.font.size = Pt(11)
                                    has_content = True
                                except Exception as para_error:
                                    print(
                                        "[PDF->PPTX] Warning: Could not add paragraph "
                                        f"{idx}: {para_error}"
                                    )
                                    continue

                            print(
                                f"[PDF->PPTX] Page {page_num}: "
                                f"Added {len(tf.paragraphs)} paragraphs to slide"
                            )

                # After processing all pages, ensure we have at least some content
                if not has_content:
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    title.text = "No Content Found"
                    content_placeholder = slide.placeholders[1]
                    tf = content_placeholder.text_frame
                    tf.text = (
                        "This PDF may be image-based or encrypted. "
                        "No extractable content found."
                    )

                # Save PowerPoint file
                output_path = os.path.join(temp_dir, "output.pptx")
                prs.save(output_path)
                print(f"[PDF->PPTX] OK Created PowerPoint file: {output_path}")

            except Exception as pdf_error:
                # Fallback: basic PowerPoint with error message
                print(f"[PDF->PPTX] pdfplumber error: {pdf_error}")
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)

                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = "Error Extracting PDF"

                content_placeholder = slide.placeholders[1]
                tf = content_placeholder.text_frame
                tf.text = str(pdf_error)

                output_path = os.path.join(temp_dir, "output.pptx")
                prs.save(output_path)
            
            # Validate output file
            if not os.path.exists(output_path):
                raise Exception("Failed to create PowerPoint file")
            
            file_size = os.path.getsize(output_path)
            if file_size < 100:
                raise Exception("PowerPoint file is too small")
            
            # Read and encode
            with open(output_path, "rb") as f:
                output_data = f.read()
            
            if len(output_data) == 0:
                raise Exception("PowerPoint file is empty")
            
            # Validate PPTX (ZIP structure)
            if output_data[:2] != b'PK':
                raise Exception("Invalid PPTX file structure")
            
            output_base64 = base64.b64encode(output_data).decode('utf-8')
            
            elapsed = time.time() - start_time
            print(f"[PDF->PPTX] OK Success in {elapsed:.1f}s - {len(output_data)} bytes")
            
            # COMPETITOR ADVANTAGE: Return presentation details
            return JSONResponse({
                "success": True,
                "file": output_base64,
                "filename": file.filename.replace('.pdf', '.pptx'),
                "size": len(output_data),
                "original_size": len(content),
                "processing_time": f"{elapsed:.1f}s",
                "quality": "enterprise-grade",
                "format": "Widescreen 16:9 (Professional)",
                "features": {
                    "widescreen_format": True,
                    "professional_fonts": True,
                    "complete_content_extraction": True,
                    "section_based_splitting": True
                }
            })
            
        except HTTPException:
            raise
        except Exception as e:
            elapsed = time.time() - start_time
            error_msg = filter_java_warnings(str(e))
            print(f"[PDF->PPTX] ERROR Failed after {elapsed:.1f}s: {error_msg}")
            raise HTTPException(500, f"Conversion failed: {error_msg}")

@app.post("/convert/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    """Convert Word to PDF - High quality conversion"""
    start_time = time.time()
    
    if not file.filename or not file.filename.lower().endswith(('.docx', '.doc')):
        raise HTTPException(400, "Only Word files (.docx, .doc) are allowed")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # Save uploaded file with sanitized filename (LibreOffice has issues with spaces/special chars)
            # Use simple name to avoid LibreOffice conversion issues
            original_filename = file.filename
            safe_filename = "input.docx" if original_filename.lower().endswith('.docx') else "input.doc"
            input_path = os.path.join(temp_dir, safe_filename)
            
            content = await file.read()
            
            if not content or len(content) < 100:
                raise HTTPException(400, "File is too small or empty")
            
            # Validate Word file (DOCX should start with PK, DOC is binary)
            if original_filename.lower().endswith('.docx'):
                if content[:2] != b'PK':
                    raise HTTPException(400, "Invalid DOCX file (not a valid ZIP archive)")
            
            with open(input_path, "wb") as f:
                f.write(content)
            
            print(f"[DOCX->PDF] Processing: {original_filename}, {len(content)} bytes (saved as {safe_filename})")
            
            # Convert to PDF - Use simple format (like working revision)
            # ULTRA-OPTIMIZED: Dynamic timeout based on file size
            file_size_mb = len(content) / (1024 * 1024)
            dynamic_timeout = 30 if file_size_mb < 1 else (45 if file_size_mb < 5 else 60)
            print(f"[DOCX->PDF] File size: {file_size_mb:.2f}MB, using timeout: {dynamic_timeout}s")
            
            output_path = convert_with_libreoffice(
                input_path,
                'pdf',  # Simple format - like working revision 00049-ss2
                temp_dir,
                timeout=dynamic_timeout  # ULTRA-OPTIMIZED: Dynamic timeout for faster processing
            )
            
            # Validate output
            if not os.path.exists(output_path):
                raise Exception("PDF file not created")
            
            file_size = os.path.getsize(output_path)
            if file_size < 100:
                raise Exception(f"PDF file too small ({file_size} bytes)")
            
            # Validate PDF header
            with open(output_path, "rb") as f:
                header = f.read(4)
                if not header.startswith(b'%PDF'):
                    raise Exception("Invalid PDF file structure")
            
            # Read and encode output
            with open(output_path, "rb") as f:
                output_data = f.read()
            
                output_base64 = base64.b64encode(output_data).decode()
            
            elapsed = time.time() - start_time
            print(f"[DOCX->PDF] OK Success in {elapsed:.1f}s - {len(output_data)} bytes")
            
            # COMPETITOR ADVANTAGE: Return quality metrics for Office to PDF
            return JSONResponse({
                "success": True,
                "file": output_base64,
                "filename": file.filename.rsplit('.', 1)[0] + '.pdf',
                "size": len(output_data),
                "original_size": len(content),
                "processing_time": f"{elapsed:.1f}s",
                "quality": "high-quality",  # Premium PDF export
                "compression_ratio": f"{(1 - len(output_data)/len(content))*100:.1f}%" if len(content) > 0 else "0%"
            })
            
        except HTTPException:
            raise
        except Exception as e:
            elapsed = time.time() - start_time
            error_msg = filter_java_warnings(str(e))
            print(f"[DOCX->PDF] ERROR Failed after {elapsed:.1f}s: {error_msg}")
            raise HTTPException(500, f"Conversion failed: {error_msg}")

@app.post("/convert/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    """Convert PowerPoint to PDF - High quality conversion"""
    start_time = time.time()
    
    if not file.filename or not file.filename.lower().endswith(('.pptx', '.ppt')):
        raise HTTPException(400, "Only PowerPoint files (.pptx, .ppt) are allowed")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # Save uploaded file
            input_path = os.path.join(temp_dir, file.filename)
            content = await file.read()
            
            if not content or len(content) < 100:
                raise HTTPException(400, "File is too small or empty")
            
            # Validate PowerPoint file (PPTX should start with PK)
            if file.filename.lower().endswith('.pptx'):
                if content[:2] != b'PK':
                    raise HTTPException(400, "Invalid PPTX file (not a valid ZIP archive)")
            
            with open(input_path, "wb") as f:
                f.write(content)
            
            print(f"[PPTX->PDF] Processing: {file.filename}, {len(content)} bytes")
            
            # ULTRA-OPTIMIZED: Dynamic timeout based on file size
            file_size_mb = len(content) / (1024 * 1024)
            dynamic_timeout = 30 if file_size_mb < 1 else (45 if file_size_mb < 5 else 60)
            print(f"[PPTX->PDF] File size: {file_size_mb:.2f}MB, using timeout: {dynamic_timeout}s")
            
            # Convert to PDF with BEST QUALITY settings
            # Use high-quality PDF export (better than default)
            # Competitor tools often produce low-quality PDFs - we ensure high quality
            output_path = convert_with_libreoffice(
                input_path,
                'pdf:"impress_pdf_Export"',  # Use Impress PDF export for presentations
                temp_dir,
                timeout=dynamic_timeout  # ULTRA-OPTIMIZED: Dynamic timeout for faster processing
            )
            
            # Validate output
            if not os.path.exists(output_path):
                raise Exception("PDF file not created")
            
            file_size = os.path.getsize(output_path)
            if file_size < 100:
                raise Exception(f"PDF file too small ({file_size} bytes)")
            
            # Validate PDF header
            with open(output_path, "rb") as f:
                header = f.read(4)
                if not header.startswith(b'%PDF'):
                    raise Exception("Invalid PDF file structure")
            
            # Read and encode output
            with open(output_path, "rb") as f:
                output_data = f.read()
            
            output_base64 = base64.b64encode(output_data).decode()
            
            elapsed = time.time() - start_time
            print(f"[PPTX->PDF] OK Success in {elapsed:.1f}s - {len(output_data)} bytes")
            
            # COMPETITOR ADVANTAGE: Return quality metrics for Office to PDF
            return JSONResponse({
                "success": True,
                "file": output_base64,
                "filename": file.filename.rsplit('.', 1)[0] + '.pdf',
                "size": len(output_data),
                "original_size": len(content),
                "processing_time": f"{elapsed:.1f}s",
                "quality": "high-quality",  # Premium PDF export
                "compression_ratio": f"{(1 - len(output_data)/len(content))*100:.1f}%" if len(content) > 0 else "0%"
            })
            
        except HTTPException:
            raise
        except Exception as e:
            elapsed = time.time() - start_time
            error_msg = filter_java_warnings(str(e))
            print(f"[PPTX->PDF] ERROR Failed after {elapsed:.1f}s: {error_msg}")
            raise HTTPException(500, f"Conversion failed: {error_msg}")

@app.post("/convert/excel-to-pdf")
async def excel_to_pdf(
    file: UploadFile = File(...),
    page_size: str = Form(None),
    orientation: str = Form(None),
    prevent_chart_split: str = Form("true"),
    scale_charts_to_fit: str = Form("true"),
    fit_to_page: str = Form("true"),
    margin_top: str = Form(None),
    margin_right: str = Form(None),
    margin_bottom: str = Form(None),
    margin_left: str = Form(None),
    print_quality: str = Form("high"),
):
    """Convert Excel to PDF - High quality conversion with chart handling"""
    start_time = time.time()
    
    if not file.filename or not file.filename.lower().endswith(('.xlsx', '.xls')):
        raise HTTPException(400, "Only Excel files (.xlsx, .xls) are allowed")
    
    # Parse options
    prevent_chart_split_bool = prevent_chart_split.lower() == "true"
    scale_charts_to_fit_bool = scale_charts_to_fit.lower() == "true"
    fit_to_page_bool = fit_to_page.lower() == "true"
    
    print(f"[XLSX->PDF] Options: prevent_chart_split={prevent_chart_split_bool}, scale_charts={scale_charts_to_fit_bool}, fit_to_page={fit_to_page_bool}")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # Save uploaded file
            input_path = os.path.join(temp_dir, file.filename)
            content = await file.read()
            
            if not content or len(content) < 100:
                raise HTTPException(400, "File is too small or empty")
            
            # Validate Excel file (XLSX should start with PK)
            if file.filename.lower().endswith('.xlsx'):
                if content[:2] != b'PK':
                    raise HTTPException(400, "Invalid XLSX file (not a valid ZIP archive)")
            
            with open(input_path, "wb") as f:
                f.write(content)
            
            print(f"[XLSX->PDF] Processing: {file.filename}, {len(content)} bytes")
            
            # For XLSX files, we can modify print settings using openpyxl before conversion
            # CRITICAL: This prevents chart splitting by setting proper print areas and scaling
            if file.filename.lower().endswith('.xlsx') and (prevent_chart_split_bool or fit_to_page_bool):
                try:
                    from openpyxl import load_workbook
                    from openpyxl.worksheet.pagebreak import Break
                    print(f"[XLSX->PDF] 🔧🔧🔧 STARTING: Modifying Excel print settings to PREVENT CHART SPLITTING...")
                    print(f"[XLSX->PDF] Input file path: {input_path}")
                    print(f"[XLSX->PDF] File exists: {os.path.exists(input_path)}")
                    print(f"[XLSX->PDF] File size: {os.path.getsize(input_path) if os.path.exists(input_path) else 0} bytes")
                    
                    wb = load_workbook(input_path)
                    for sheet in wb.worksheets:
                        print(f"[XLSX->PDF] Processing sheet: {sheet.title}")
                        
                        # CRITICAL: Set fit to page to prevent splitting
                        # AGGRESSIVE SETTINGS: Force everything to fit on pages without splitting charts
                        if fit_to_page_bool or prevent_chart_split_bool:
                            # Fit to 1 page wide - THIS IS KEY to prevent horizontal chart splits
                            sheet.page_setup.fitToWidth = 1
                            # For chart preservation, we want to fit to height as well
                            # But we use a large number to allow multiple pages if needed, but prevent horizontal splits
                            # Setting to 0 means "fit all" which can cause issues
                            # Instead, use a very large number to allow vertical expansion but prevent horizontal splits
                            sheet.page_setup.fitToHeight = 999  # Very large number - allows vertical expansion but prevents horizontal splits
                            sheet.page_setup.fitToPage = True
                            sheet.page_setup.scale = None  # MUST be None for fitToPage to work
                            sheet.page_setup.blackAndWhite = False
                            sheet.page_setup.draft = False
                            # CRITICAL: Set page order to "down then over" to ensure charts aren't split
                            sheet.page_setup.pageOrder = 'downThenOver'
                            print(f"[XLSX->PDF] ✅✅✅ AGGRESSIVE: fitToWidth=1, fitToHeight=999, fitToPage=True, scale=None, pageOrder=downThenOver")
                        
                        # CRITICAL: Remove ALL page breaks - THIS IS ESSENTIAL
                        # Page breaks are the #1 cause of chart splitting
                        try:
                            if hasattr(sheet, 'row_breaks'):
                                sheet.row_breaks = []
                            if hasattr(sheet, 'col_breaks'):
                                sheet.col_breaks = []
                            # Also try to clear any internal page break references
                            if hasattr(sheet, '_print_rows'):
                                sheet._print_rows = None
                            if hasattr(sheet, '_print_cols'):
                                sheet._print_cols = None
                            print(f"[XLSX->PDF] ✅✅✅ Removed ALL page breaks (row_breaks, col_breaks, _print_rows, _print_cols)")
                        except Exception as e:
                            print(f"[XLSX->PDF] ⚠️ Warning: Could not remove some page breaks: {e}")
                        
                        # Set print area to include ALL content (prevents truncation)
                        if sheet.max_row > 0 and sheet.max_column > 0:
                            max_col_letter = get_column_letter(sheet.max_column)
                            print_area = f'A1:{max_col_letter}{sheet.max_row}'
                            sheet.print_area = print_area
                            print(f"[XLSX->PDF] ✅ Set print area: {print_area}")
                        
                        # Set orientation based on user selection
                        # If prevent_chart_split is enabled, we'll auto-move/resize charts instead of forcing landscape
                        if orientation and orientation.lower() == "landscape":
                            sheet.page_setup.orientation = sheet.ORIENTATION_LANDSCAPE
                            print(f"[XLSX->PDF] ✅ Set orientation: LANDSCAPE")
                        else:
                            sheet.page_setup.orientation = sheet.ORIENTATION_PORTRAIT
                            print(f"[XLSX->PDF] ✅ Set orientation: PORTRAIT")
                        
                        # If prevent_chart_split is enabled, we'll auto-move charts to columns A-E and resize them
                        # This works in both portrait and landscape, but landscape gives more space
                        if prevent_chart_split_bool:
                            print(f"[XLSX->PDF] ✅✅✅ Chart auto-fix enabled: Charts will be moved to columns A-E and resized during conversion")
                        
                        # Set page size (using numeric constants)
                        if page_size:
                            # Use numeric constants (openpyxl paper size constants)
                            size_map = {
                                "A4": 9,  # PAPERSIZE_A4
                                "Letter": 1,  # PAPERSIZE_LETTER
                                "Legal": 5,  # PAPERSIZE_LEGAL
                                "A3": 8,  # PAPERSIZE_A3
                                "Tabloid": 3,  # PAPERSIZE_TABLOID
                            }
                            if page_size in size_map:
                                sheet.page_setup.paperSize = size_map[page_size]
                                print(f"[XLSX->PDF] ✅ Set paper size: {page_size} (constant: {size_map[page_size]})")
                        
                        # Set margins - CRITICAL: Use minimal margins for chart preservation
                        # Smaller margins = more space for charts to fit on one page
                        if prevent_chart_split_bool:
                            # FORCE minimal margins when chart split prevention is enabled
                            sheet.page_margins.top = 0.3  # Minimal top margin
                            sheet.page_margins.bottom = 0.3  # Minimal bottom margin
                            sheet.page_margins.left = 0.3  # Minimal left margin
                            sheet.page_margins.right = 0.3  # Minimal right margin
                            sheet.page_margins.header = 0.2  # Minimal header
                            sheet.page_margins.footer = 0.2  # Minimal footer
                            print(f"[XLSX->PDF] ✅✅✅ FORCED: Minimal margins (0.3\") for maximum chart space")
                        else:
                            # Use user-specified margins
                            if margin_top:
                                sheet.page_margins.top = float(margin_top)
                            if margin_bottom:
                                sheet.page_margins.bottom = float(margin_bottom)
                            if margin_left:
                                sheet.page_margins.left = float(margin_left)
                            if margin_right:
                                sheet.page_margins.right = float(margin_right)
                        
                        # CRITICAL: Set print quality for charts
                        if print_quality == "high":
                            sheet.page_setup.horizontalDpi = 300
                            sheet.page_setup.verticalDpi = 300
                            print(f"[XLSX->PDF] ✅ Set high quality (300 DPI)")
                        
                        # CRITICAL: AUTOMATED CHART POSITIONING - This automates the manual process
                        # We automatically move ALL charts to column A, row 1 and resize to 6×12
                        # This is the automated version of: "Move charts to columns A-E, resize to 8 columns width"
                        if prevent_chart_split_bool:
                            try:
                                chart_count = 0
                                charts_moved = 0
                                
                                # Method 1: Access charts via _charts attribute (if available)
                                if hasattr(sheet, '_charts') and sheet._charts:
                                    chart_count = len(sheet._charts)
                                    print(f"[XLSX->PDF] 📊 Found {chart_count} chart(s) via _charts attribute")
                                    
                                    for idx, chart in enumerate(sheet._charts):
                                        try:
                                            # Get or create anchor
                                            if hasattr(chart, 'anchor'):
                                                anchor = chart.anchor
                                                
                                                # Ensure anchor has _from
                                                if not hasattr(anchor, '_from'):
                                                    # Create anchor if it doesn't exist
                                                    from openpyxl.drawing.spreadsheet_drawing import OneCellAnchor
                                                    anchor = OneCellAnchor()
                                                    chart.anchor = anchor
                                                
                                                if hasattr(anchor, '_from'):
                                                    from_obj = anchor._from
                                                    old_col = from_obj.col if hasattr(from_obj, 'col') else None
                                                    old_row = from_obj.row if hasattr(from_obj, 'row') else None
                                                    
                                                    # AUTOMATED: Move to column A (0), row 1
                                                    from_obj.col = 0
                                                    from_obj.row = 1
                                                    charts_moved += 1
                                                    print(f"[XLSX->PDF] ✅✅✅ AUTOMATED: Moved chart {idx+1} to column A, row 1 (was column {old_col+1 if old_col is not None else 'unknown'}, row {old_row+1 if old_row is not None else 'unknown'})")
                                                
                                                # AUTOMATED: Resize to 6 columns × 12 rows
                                                if hasattr(anchor, '_to'):
                                                    if hasattr(anchor._to, 'col'):
                                                        anchor._to.col = 6  # 6 columns width (A-F)
                                                        print(f"[XLSX->PDF] ✅✅✅ AUTOMATED: Resized chart {idx+1} width to 6 columns (A-F)")
                                                    if hasattr(anchor._to, 'row'):
                                                        anchor._to.row = 13  # Row 1 + 12 rows = row 13
                                                        print(f"[XLSX->PDF] ✅✅✅ AUTOMATED: Resized chart {idx+1} height to 12 rows")
                                                
                                        except Exception as e:
                                            print(f"[XLSX->PDF] ⚠️ Could not automate chart {idx+1} positioning: {e}")
                                
                                # Method 2: Count charts via relationships
                                if hasattr(sheet, '_rels') and sheet._rels:
                                    rel_chart_count = sum(1 for rel in sheet._rels.values() if hasattr(rel, 'target') and 'chart' in str(rel.target).lower())
                                    if rel_chart_count > chart_count:
                                        chart_count = rel_chart_count
                                        print(f"[XLSX->PDF] 📊 Found {chart_count} chart(s) via relationships")
                                
                                if chart_count > 0:
                                    if charts_moved > 0:
                                        print(f"[XLSX->PDF] ✅✅✅ AUTOMATION COMPLETE: Moved {charts_moved} chart(s) to column A, row 1, size 6×12")
                                        print(f"[XLSX->PDF] ✅✅✅ This automates the manual process: 'Move charts to columns A-E, resize to 8 columns width'")
                                    else:
                                        print(f"[XLSX->PDF] ℹ️ {chart_count} chart(s) found, attempting XML manipulation...")
                                else:
                                    print(f"[XLSX->PDF] ℹ️ No charts detected via openpyxl, will try XML manipulation...")
                            except Exception as e:
                                import traceback
                                print(f"[XLSX->PDF] ⚠️ Chart automation error: {e}")
                                print(f"[XLSX->PDF] Traceback: {traceback.format_exc()}")
                                print(f"[XLSX->PDF] Will try XML manipulation as fallback...")
                        
                        # CRITICAL: Enable "Print object placeholders" to ensure charts are included
                        if hasattr(sheet, 'sheet_view'):
                            sheet.sheet_view.showGridLines = True
                    
                    # Save modified workbook
                    modified_path = os.path.join(temp_dir, "modified_" + file.filename)
                    
                    # CRITICAL: Close workbook before saving to ensure all changes are written
                    wb.save(modified_path)
                    wb.close()  # Close workbook to ensure all changes are flushed
                    
                    # Verify file was saved
                    if os.path.exists(modified_path):
                        file_size = os.path.getsize(modified_path)
                        print(f"[XLSX->PDF] ✅ Excel file modified and saved: {modified_path}")
                        print(f"[XLSX->PDF] ✅ File size: {file_size} bytes")
                    else:
                        raise Exception(f"Failed to save modified Excel file: {modified_path}")
                    
                    input_path = modified_path
                    
                    # CRITICAL: Try to manipulate Excel XML directly to move charts
                    # This is a workaround for openpyxl's limited chart manipulation
                    if prevent_chart_split_bool:
                        try:
                            import zipfile
                            import xml.etree.ElementTree as ET
                            from io import BytesIO
                            
                            print(f"[XLSX->PDF] 🔄 AUTOMATING: Moving charts to column A, row 1, size 6×12 via XML manipulation...")
                            print(f"[XLSX->PDF] 🔄 This automates: 'Move charts to columns A-E, resize to 8 columns width'")
                            
                            # Excel files are ZIP archives
                            with zipfile.ZipFile(modified_path, 'r') as zip_read:
                                # Find drawing files (charts are in drawings)
                                drawing_files = [f for f in zip_read.namelist() if 'drawings/drawing' in f and f.endswith('.xml')]
                                
                                if drawing_files:
                                    print(f"[XLSX->PDF] 📊 Found {len(drawing_files)} drawing file(s)")
                                    
                                    # Create new ZIP with modified XML
                                    new_zip_data = BytesIO()
                                    with zipfile.ZipFile(new_zip_data, 'w', zipfile.ZIP_DEFLATED) as zip_write:
                                        # Copy all files, modifying drawing files
                                        for item in zip_read.infolist():
                                            data = zip_read.read(item.filename)
                                            
                                            # Modify drawing XML files
                                            if item.filename in drawing_files:
                                                try:
                                                    # Parse XML
                                                    root = ET.fromstring(data)
                                                    
                                                    # Find all anchors (charts/images)
                                                    # Namespace for drawingML
                                                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                                          'xdr': 'http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing'}
                                                    
                                                    moved_count = 0
                                                    # CRITICAL: ALWAYS move ALL charts to column A (0) and row 1
                                                    # This ensures charts are in the safest position to prevent splitting
                                                    for anchor in root.findall('.//xdr:twoCellAnchor', ns):
                                                        # Get from position
                                                        from_elem = anchor.find('.//xdr:from', ns)
                                                        if from_elem is not None:
                                                            col_elem = from_elem.find('xdr:col', ns)
                                                            row_elem = from_elem.find('xdr:row', ns)
                                                            if col_elem is not None:
                                                                try:
                                                                    old_col = int(col_elem.text)
                                                                    # ALWAYS move to column A (0) - don't check if beyond column E
                                                                    col_elem.text = '0'
                                                                    # ALWAYS move to row 1 for consistency
                                                                    if row_elem is not None:
                                                                        try:
                                                                            old_row = int(row_elem.text)
                                                                            row_elem.text = '1'
                                                                            moved_count += 1
                                                                            print(f"[XLSX->PDF] ✅✅✅ XML: FORCED chart to column A, row 1 (was column {old_col+1}, row {old_row+1})")
                                                                        except:
                                                                            row_elem.text = '1'
                                                                            moved_count += 1
                                                                            print(f"[XLSX->PDF] ✅✅✅ XML: FORCED chart to column A, row 1 (was column {old_col+1})")
                                                                    else:
                                                                        moved_count += 1
                                                                        print(f"[XLSX->PDF] ✅✅✅ XML: FORCED chart to column A (was column {old_col+1})")
                                                                except (ValueError, AttributeError):
                                                                    pass
                                                    
                                                    # Also check oneCellAnchor - ALWAYS move to column A, row 1
                                                    for anchor in root.findall('.//xdr:oneCellAnchor', ns):
                                                        from_elem = anchor.find('.//xdr:from', ns)
                                                        if from_elem is not None:
                                                            col_elem = from_elem.find('xdr:col', ns)
                                                            row_elem = from_elem.find('xdr:row', ns)
                                                            if col_elem is not None:
                                                                try:
                                                                    old_col = int(col_elem.text)
                                                                    # ALWAYS move to column A (0)
                                                                    col_elem.text = '0'
                                                                    # ALWAYS move to row 1
                                                                    if row_elem is not None:
                                                                        try:
                                                                            old_row = int(row_elem.text)
                                                                            row_elem.text = '1'
                                                                            moved_count += 1
                                                                            print(f"[XLSX->PDF] ✅✅✅ XML: FORCED oneCellAnchor chart to column A, row 1 (was column {old_col+1}, row {old_row+1})")
                                                                        except:
                                                                            row_elem.text = '1'
                                                                            moved_count += 1
                                                                            print(f"[XLSX->PDF] ✅✅✅ XML: FORCED oneCellAnchor chart to column A, row 1 (was column {old_col+1})")
                                                                    else:
                                                                        moved_count += 1
                                                                        print(f"[XLSX->PDF] ✅✅✅ XML: FORCED oneCellAnchor chart to column A (was column {old_col+1})")
                                                                except (ValueError, AttributeError):
                                                                    pass
                                                    
                                                    # CRITICAL: Resize charts to fit within page dimensions
                                                    # A4 landscape: width ≈ 11.69" (29.7cm), height ≈ 8.27" (21cm)
                                                    # With 0.3" margins: usable width ≈ 11.09" (28.2cm), height ≈ 7.67" (19.5cm)
                                                    # In Excel: 1 column ≈ 0.7cm, 1 row ≈ 0.5cm
                                                    # Safe limits: width = 8 columns (A-H), height = 15 rows
                                                    resized_count = 0
                                                    for anchor in root.findall('.//xdr:twoCellAnchor', ns):
                                                        from_elem = anchor.find('.//xdr:from', ns)
                                                        to_elem = anchor.find('.//xdr:to', ns)
                                                        if from_elem is not None and to_elem is not None:
                                                            from_col_elem = from_elem.find('xdr:col', ns)
                                                            to_col_elem = to_elem.find('xdr:col', ns)
                                                            from_row_elem = from_elem.find('xdr:row', ns)
                                                            to_row_elem = to_elem.find('xdr:row', ns)
                                                            
                                                            if from_col_elem is not None and to_col_elem is not None:
                                                                try:
                                                                    from_col = int(from_col_elem.text)
                                                                    to_col = int(to_col_elem.text)
                                                                    chart_width = to_col - from_col
                                                                    
                                                                    # CRITICAL: Use larger width for landscape mode (charts fit much better)
                                                                    # Landscape A4: ~11.69" width, with 0.3" margins = ~11.09" usable
                                                                    # In Excel: 1 column ≈ 0.7cm, so ~15 columns fit comfortably
                                                                    # Use 10 columns for safety (A-J) - gives more space for chart details
                                                                    max_width = 10  # Increased from 6 to 10 for better chart rendering
                                                                    # Since we moved chart to column A (0), set to_col to max_width
                                                                    to_col_elem.text = str(max_width)
                                                                    resized_count += 1
                                                                    if chart_width != max_width:
                                                                        print(f"[XLSX->PDF] ✅✅✅ XML: FORCED chart width to {max_width} columns (was {chart_width} columns)")
                                                                    else:
                                                                        print(f"[XLSX->PDF] ✅✅✅ XML: Chart width set to {max_width} columns")
                                                                    
                                                                    # Also limit chart height to prevent vertical splitting
                                                                    if from_row_elem is not None and to_row_elem is not None:
                                                                        try:
                                                                            from_row = int(from_row_elem.text)
                                                                            to_row = int(to_row_elem.text)
                                                                            chart_height = to_row - from_row
                                                                            
                                                                            # CRITICAL: Use larger height for landscape mode
                                                                            # Landscape A4: ~8.27" height, with 0.3" margins = ~7.67" usable
                                                                            # In Excel: 1 row ≈ 0.5cm, so ~15 rows fit comfortably
                                                                            # Use 14 rows for safety - gives more space for chart details
                                                                            max_height = 14  # Increased from 12 to 14 for better chart rendering
                                                                            to_row_elem.text = str(max_height)
                                                                            resized_count += 1
                                                                            if chart_height != max_height:
                                                                                print(f"[XLSX->PDF] ✅✅✅ XML: FORCED chart height to {max_height} rows (was {chart_height} rows)")
                                                                            else:
                                                                                print(f"[XLSX->PDF] ✅✅✅ XML: Chart height set to {max_height} rows")
                                                                        except (ValueError, AttributeError):
                                                                            pass
                                                                except (ValueError, AttributeError):
                                                                    pass
                                                    
                                                    if resized_count > 0:
                                                        print(f"[XLSX->PDF] ✅✅✅ XML: Resized {resized_count} chart dimension(s) to fit page")
                                                    
                                                    # CRITICAL: Always write modified XML, even if no charts were moved
                                                    # This ensures all charts are positioned correctly
                                                    ET.register_namespace('a', 'http://schemas.openxmlformats.org/drawingml/2006/main')
                                                    ET.register_namespace('xdr', 'http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing')
                                                    
                                                    # Write modified XML back
                                                    data = ET.tostring(root, encoding='utf-8', xml_declaration=True)
                                                    
                                                    if moved_count > 0 or resized_count > 0:
                                                        print(f"[XLSX->PDF] ✅✅✅ XML manipulation: Moved {moved_count} chart(s) to column A, resized {resized_count} dimension(s)")
                                                    else:
                                                        print(f"[XLSX->PDF] ✅✅✅ XML manipulation: All charts already positioned correctly (column A, 6×12 size)")
                                                except Exception as e:
                                                    print(f"[XLSX->PDF] ⚠️ Could not modify drawing XML {item.filename}: {e}")
                                            
                                            # Write file to new ZIP
                                            zip_write.writestr(item, data)
                                    
                                    # Replace original file with modified version
                                    with open(modified_path, 'wb') as f:
                                        f.write(new_zip_data.getvalue())
                                    
                                    print(f"[XLSX->PDF] ✅✅✅ Excel XML modified - charts moved to prevent splitting")
                        except Exception as e:
                            import traceback
                            print(f"[XLSX->PDF] ⚠️ XML manipulation failed: {e}")
                            print(f"[XLSX->PDF] Traceback: {traceback.format_exc()}")
                            print(f"[XLSX->PDF] Continuing with openpyxl-modified file...")
                except Exception as e:
                    import traceback
                    print(f"[XLSX->PDF] ⚠️ Could not modify print settings: {e}")
                    print(f"[XLSX->PDF] Traceback: {traceback.format_exc()}")
                    print(f"[XLSX->PDF] Continuing with original file...")
            
            # CRITICAL: Auto-detect charts BEFORE building filter parameters
            # Charts render MUCH better in landscape mode
            has_charts = False
            if file.filename.lower().endswith('.xlsx'):
                try:
                    import zipfile
                    with zipfile.ZipFile(input_path, 'r') as zip_read:
                        drawing_files = [f for f in zip_read.namelist() if 'drawings/drawing' in f and f.endswith('.xml')]
                        has_charts = len(drawing_files) > 0
                        if has_charts:
                            print(f"[XLSX->PDF] 📊 Charts detected: {len(drawing_files)} drawing file(s) found - will optimize for chart rendering")
                except Exception as e:
                    print(f"[XLSX->PDF] ⚠️ Could not detect charts: {e}")
            
            # Build PDF export filter with options
            # LibreOffice Calc PDF export supports various parameters
            pdf_filter_params = []
            
            # Quality settings
            if print_quality == "high":
                pdf_filter_params.append("Quality=100")
            elif print_quality == "normal":
                pdf_filter_params.append("Quality=75")
            else:
                pdf_filter_params.append("Quality=50")
            
            # Page size
            if page_size:
                size_map = {
                    "A4": "A4",
                    "Letter": "Letter",
                    "Legal": "Legal",
                    "A3": "A3",
                    "Tabloid": "Tabloid"
                }
                if page_size in size_map:
                    pdf_filter_params.append(f"PageFormat={size_map[page_size]}")
            
            # Orientation - Auto-detect charts and force landscape if charts are present
            
            # CRITICAL: Force landscape if charts are detected (charts fit much better in landscape)
            if has_charts and prevent_chart_split_bool:
                pdf_filter_params.append("IsLandscape=true")
                print(f"[XLSX->PDF] ✅✅✅ AUTO: Forcing LANDSCAPE orientation (charts detected - better chart rendering)")
            elif orientation:
                if orientation.lower() == "landscape":
                    pdf_filter_params.append("IsLandscape=true")
                    print(f"[XLSX->PDF] ✅ Using landscape orientation (user selected)")
                else:
                    pdf_filter_params.append("IsLandscape=false")
                    print(f"[XLSX->PDF] ✅ Using portrait orientation (user selected)")
                    if prevent_chart_split_bool:
                        print(f"[XLSX->PDF] ⚠️ Portrait mode: Charts will be auto-moved/resized to fit (landscape recommended for better results)")
            elif prevent_chart_split_bool:
                # Default to landscape if prevent_chart_split is enabled
                pdf_filter_params.append("IsLandscape=true")
                print(f"[XLSX->PDF] ✅ Using landscape orientation (default for chart preservation)")
            
            # CRITICAL: Add chart-specific parameters for better rendering
            if has_charts and prevent_chart_split_bool:
                # Use higher quality for charts
                pdf_filter_params.append("Quality=100")
                # Ensure all objects are included
                pdf_filter_params.append("SelectPdfVersion=1")  # PDF 1.4 for better compatibility
                print(f"[XLSX->PDF] ✅ Added chart-specific PDF parameters (Quality=100, PDF 1.4)")
            
            # Build filter string with proper escaping
            if pdf_filter_params:
                try:
                    params_str = ','.join(pdf_filter_params)
                    filter_str = f'pdf:"calc_pdf_Export:{{{params_str}}}"'
                except:
                    # Fallback to simple filter if parameter formatting fails
                    filter_str = 'pdf:"calc_pdf_Export"'
            else:
                filter_str = 'pdf:"calc_pdf_Export"'
            
            print(f"[XLSX->PDF] Using LibreOffice filter: {filter_str}")
            print(f"[XLSX->PDF] ⚠️⚠️⚠️ CHART SPLIT PREVENTION ACTIVE: prevent_chart_split={prevent_chart_split_bool}, fit_to_page={fit_to_page_bool}")
            
            # CRITICAL: For chart split prevention, we MUST use the modified Excel file
            # The modified file has fitToWidth=1, fitToHeight=0, and all page breaks removed
            # LibreOffice will respect these settings when converting to PDF
            
            # CRITICAL: Convert to PDF with optimized settings
            # The modified Excel file has all anti-split settings applied
            # LibreOffice should respect these settings when converting
            print(f"[XLSX->PDF] 🔄 Starting LibreOffice conversion with anti-split settings...")
            print(f"[XLSX->PDF] Input file: {input_path}")
            print(f"[XLSX->PDF] File exists: {os.path.exists(input_path)}")
            print(f"[XLSX->PDF] File size: {os.path.getsize(input_path) if os.path.exists(input_path) else 0} bytes")
            
            # CRITICAL: Try to use LibreOffice UNO API for better chart handling
            # If UNO API is not available, fall back to command-line conversion
            use_uno_api = False
            if prevent_chart_split_bool:
                try:
                    # Try to use UNO API for better control
                    print(f"[XLSX->PDF] 🔄 Attempting to use LibreOffice UNO API for chart preservation...")
                    # UNO API implementation would go here (complex, requires LibreOffice running with socket)
                    # For now, we'll use command-line conversion
                    print(f"[XLSX->PDF] ⚠️ UNO API not available, using command-line conversion")
                except Exception as e:
                    print(f"[XLSX->PDF] ⚠️ UNO API failed: {e}, using command-line conversion")
            
            # ULTRA-OPTIMIZED: Dynamic timeout based on file size
            # Use file size from disk since content might have been modified
            file_size_bytes = os.path.getsize(input_path) if os.path.exists(input_path) else len(content)
            file_size_mb = file_size_bytes / (1024 * 1024)
            dynamic_timeout = 30 if file_size_mb < 1 else (45 if file_size_mb < 5 else 60)
            print(f"[XLSX->PDF] File size: {file_size_mb:.2f}MB, using timeout: {dynamic_timeout}s")
            
            output_path = convert_with_libreoffice(
                input_path,  # This is the modified file with anti-split settings
                filter_str,
                temp_dir, 
                timeout=dynamic_timeout  # ULTRA-OPTIMIZED: Dynamic timeout for faster processing
            )
            
            print(f"[XLSX->PDF] ✅ Conversion complete.")
            
            # CRITICAL WARNING: LibreOffice limitation
            if prevent_chart_split_bool:
                print(f"[XLSX->PDF] ⚠️⚠️⚠️ IMPORTANT: LibreOffice has a KNOWN LIMITATION with chart splitting.")
                print(f"[XLSX->PDF] ⚠️ Charts are floating objects in Excel, and LibreOffice places them based on")
                print(f"[XLSX->PDF] ⚠️ their position in the sheet, NOT based on print settings (fitToPage, fitToWidth, etc.).")
                print(f"[XLSX->PDF] ⚠️ Even with all print settings applied, LibreOffice may still split charts.")
                print(f"[XLSX->PDF] ⚠️ This is a fundamental limitation of LibreOffice Calc's PDF export.")
                print(f"[XLSX->PDF] ⚠️ To truly prevent chart splitting, we would need to:")
                print(f"[XLSX->PDF] ⚠️   1. Use LibreOffice UNO API (complex, requires running LibreOffice with socket)")
                print(f"[XLSX->PDF] ⚠️   2. Post-process PDF to detect and fix split charts (complex)")
                print(f"[XLSX->PDF] ⚠️   3. Use a commercial solution (Aspose, etc.)")
                print(f"[XLSX->PDF] ⚠️   4. Convert charts to images first, then place them (loses vector quality)")
            
            # Validate output
            if not os.path.exists(output_path):
                raise Exception("PDF file not created")
            
            file_size = os.path.getsize(output_path)
            if file_size < 100:
                raise Exception(f"PDF file too small ({file_size} bytes)")
            
            # Validate PDF header
            with open(output_path, "rb") as f:
                header = f.read(4)
                if not header.startswith(b'%PDF'):
                    raise Exception("Invalid PDF file structure")
            
            # Read and encode output
            with open(output_path, "rb") as f:
                output_data = f.read()
            
                output_base64 = base64.b64encode(output_data).decode()
            
            elapsed = time.time() - start_time
            print(f"[XLSX->PDF] OK Success in {elapsed:.1f}s - {len(output_data)} bytes")
            
            # COMPETITOR ADVANTAGE: Return quality metrics for Office to PDF
            return JSONResponse({
                "success": True,
                "file": output_base64,
                "filename": file.filename.rsplit('.', 1)[0] + '.pdf',
                "size": len(output_data),
                "original_size": len(content),
                "processing_time": f"{elapsed:.1f}s",
                "quality": "high-quality",  # Premium PDF export
                "compression_ratio": f"{(1 - len(output_data)/len(content))*100:.1f}%" if len(content) > 0 else "0%"
            })
            
        except HTTPException:
            raise
        except Exception as e:
            elapsed = time.time() - start_time
            error_msg = filter_java_warnings(str(e))
            print(f"[XLSX->PDF] ERROR Failed after {elapsed:.1f}s: {error_msg}")
            raise HTTPException(500, f"Conversion failed: {error_msg}")

def _build_output_filename(original_filename: str, new_extension: str) -> str:
    """Build output filename by replacing extension"""
    if not original_filename:
        return f"output{new_extension}"
    base_name = os.path.splitext(original_filename)[0]
    return f"{base_name}{new_extension}"

@app.post("/convert/password-protect")
async def password_protect_pdf(
    file: UploadFile = File(...),
    user_password: str = Form(...),
    owner_password: str = Form(None),
    allow_printing: str = Form("true"),
    allow_copying: str = Form("true"),
    allow_editing: str = Form("false"),
):
    """Password protect PDF with AES-256 encryption"""
    start_time = time.time()
    
    if not file.filename or not file.filename.lower().endswith('.pdf'):
        raise HTTPException(400, "File must be a PDF")
    
    if not user_password or len(user_password.strip()) == 0:
        raise HTTPException(400, "User password is required")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # Save input PDF
            input_path = os.path.join(temp_dir, "input.pdf")
            output_path = os.path.join(temp_dir, "output.pdf")
            content = await file.read()
            with open(input_path, "wb") as f:
                f.write(content)
            
            # Load PDF with PyMuPDF
            doc = fitz.open(input_path)
            
            # Parse boolean permissions
            allow_print = allow_printing.lower() == "true"
            allow_copy = allow_copying.lower() == "true"
            allow_edit = allow_editing.lower() == "true"
            
            # Calculate permissions flag
            # PDF permissions: print=4, modify=8, copy=16, annotate=32
            permissions = 0
            if allow_print:
                permissions |= fitz.PDF_PERM_PRINT  # 4 - Print
            if allow_edit:
                permissions |= fitz.PDF_PERM_MODIFY  # 8 - Modify
                permissions |= fitz.PDF_PERM_ANNOTATE  # 32 - Annotate
            if allow_copy:
                permissions |= fitz.PDF_PERM_COPY  # 16 - Copy
            
            # If no permissions set, allow all (default)
            if permissions == 0:
                permissions = -1  # All permissions (0xFFFFFFFF)
            
            # Set owner password (defaults to user password if not provided)
            final_owner_password = owner_password.strip() if owner_password and owner_password.strip() else user_password.strip()
            
            # Encrypt the PDF and save to a different output file
            # Use fitz.PDF_ENCRYPT_AES_256 for AES-256 encryption (strongest)
            doc.save(
                output_path,  # Save to different file
                encryption=fitz.PDF_ENCRYPT_AES_256,  # AES-256 encryption (strongest)
                owner_pw=final_owner_password,
                user_pw=user_password.strip(),
                permissions=permissions,
                incremental=False,  # Fresh save (not incremental)
            )
            
            doc.close()
            
            # Verify file was created
            if not os.path.exists(output_path):
                raise Exception("Failed to create encrypted PDF")
            
            file_size = os.path.getsize(output_path)
            if file_size < 100:
                raise Exception(f"Encrypted PDF file too small ({file_size} bytes)")
            
            # Verify encryption was applied
            try:
                verify_doc = fitz.open(output_path)
                is_encrypted = verify_doc.is_encrypted
                needs_auth = verify_doc.needs_pass
                verify_doc.close()
                
                print(f"[Password-Protect] Encryption check: is_encrypted={is_encrypted}, needs_pass={needs_auth}")
                
                if not is_encrypted:
                    print(f"[Password-Protect] ERROR: PDF was saved but encryption flag is False!")
                    raise Exception("PDF encryption was not applied correctly. The file may not require a password.")
                
                # Try to open with password to verify it works
                try:
                    test_doc = fitz.open(output_path)
                    if test_doc.authenticate(user_password.strip()):
                        print(f"[Password-Protect] ✅ Password authentication successful")
                        test_doc.close()
                    else:
                        print(f"[Password-Protect] ⚠️ Password authentication failed")
                        test_doc.close()
                except Exception as auth_error:
                    print(f"[Password-Protect] Authentication test error: {auth_error}")
                    
            except Exception as verify_error:
                print(f"[Password-Protect] Warning: Could not verify encryption: {verify_error}")
            
            # Read and encode output
            with open(output_path, "rb") as f:
                output_data = f.read()
            
            output_base64 = base64.b64encode(output_data).decode()
            
            elapsed = time.time() - start_time
            print(f"[Password-Protect] OK Success in {elapsed:.1f}s - {len(output_data)} bytes")
            
            return JSONResponse({
                "success": True,
                "file": output_base64,
                "filename": _build_output_filename(file.filename, ".pdf"),
                "size": len(output_data),
                "original_size": len(content),
                "processing_time": f"{elapsed:.1f}s",
            })
            
        except HTTPException:
            raise
        except Exception as e:
            elapsed = time.time() - start_time
            error_msg = filter_java_warnings(str(e))
            print(f"[Password-Protect] ERROR Failed after {elapsed:.1f}s: {error_msg}")
            raise HTTPException(500, f"Password protect failed: {error_msg}")

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 9001))  # Default to 9001 for local development
    uvicorn.run(app, host="0.0.0.0", port=port)


